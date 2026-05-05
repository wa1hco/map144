#!/home/jeff/ham/map144/.venv/bin/python3
"""sync_pptx.py — two-way sync between map144_presentation.pptx and make_slides.py.

Usage (run from project root):

    python docs/sync_pptx.py           # show what would change, dry-run
    python docs/sync_pptx.py --apply   # apply all changes to make_slides.py
                                       # and rewrite the pptx without AI: lines

What it does
------------
1.  Copies docs/map144_presentation.pptx to /tmp/edited.pptx  (safety copy).
2.  Regenerates the baseline by running make_slides.py.
3.  Diffs every slide's visible text between the baseline and /tmp/edited.pptx.
    Any changed bullet / title text is patched into make_slides.py.
4.  Scans every slide's notes for lines starting with "AI:".
    Each instruction is printed (dry-run) or handed to the AI for action.
    After --apply the AI: lines are stripped from the notes in make_slides.py.

AI: instruction convention
---------------------------
In LibreOffice Impress, click the notes panel at the bottom of any slide and
add one or more lines starting with  AI:  followed by a plain-English request:

    AI: define what SPD stands for in the second bullet
    AI: split this slide into two — detection on one, decode on the other
    AI: expand the last bullet with one sentence of explanation

Everything else in the notes panel is treated as speaker notes and left alone.
"""

import argparse
import re
import shutil
import subprocess
import sys
from pathlib import Path

from pptx import Presentation

PPTX     = Path("docs/map144_presentation.pptx")
SCRIPT   = Path("docs/make_slides.py")
EDITED   = Path("/tmp/map144_edited.pptx")
BASELINE = Path("/tmp/map144_baseline.pptx")


# ── PPTX text extraction ──────────────────────────────────────────────────────

def extract_slides(path: Path) -> list[dict]:
    """Return list of {slide_num, title, texts, notes, ai_notes} dicts."""
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            texts.append(line)
            except Exception:
                pass

        notes_text = ""
        if slide.has_notes_slide:
            try:
                notes_text = slide.notes_slide.notes_text_frame.text
            except Exception:
                pass

        note_lines   = notes_text.splitlines()
        ai_lines     = [l[2:].strip() for l in note_lines if l.strip().startswith("AI:")]
        clean_notes  = "\n".join(l for l in note_lines if not l.strip().startswith("AI:")).strip()

        title = texts[0] if texts else f"(slide {i+1})"
        slides.append({
            "slide_num":  i + 1,
            "title":      title,
            "texts":      texts,
            "notes":      clean_notes,
            "ai_notes":   ai_lines,
        })
    return slides


# ── Text diff ─────────────────────────────────────────────────────────────────

def find_text_diffs(baseline: list[dict], edited: list[dict]) -> list[dict]:
    """Return slides where visible text changed."""
    diffs = []
    for b, e in zip(baseline, edited):
        removed = [t for t in b["texts"] if t not in e["texts"]]
        added   = [t for t in e["texts"] if t not in b["texts"]]
        if removed or added:
            diffs.append({
                "slide_num": b["slide_num"],
                "title":     b["title"],
                "removed":   removed,
                "added":     added,
            })
    return diffs


# ── make_slides.py patching ───────────────────────────────────────────────────

def patch_script(diffs: list[dict]) -> list[str]:
    """Apply text diffs to make_slides.py.  Returns list of applied patches."""
    source = SCRIPT.read_text()
    applied = []
    skipped = []

    for diff in diffs:
        # Only handle simple 1-for-1 replacements; flag anything else for manual review
        if len(diff["removed"]) != 1 or len(diff["added"]) != 1:
            skipped.append(diff)
            continue

        old_text = diff["removed"][0]
        new_text = diff["added"][0]

        # Escape for use as a literal string search
        if old_text not in source:
            skipped.append(diff)
            continue

        source = source.replace(old_text, new_text, 1)
        applied.append(f"Slide {diff['slide_num']} ({diff['title']!r}): "
                       f"{old_text!r} → {new_text!r}")

    SCRIPT.write_text(source)

    if skipped:
        print("\n⚠  Could not auto-patch these diffs (not a simple 1-for-1 replacement):")
        for d in skipped:
            print(f"  Slide {d['slide_num']} ({d['title']!r})")
            for r in d["removed"]: print(f"    - {r!r}")
            for a in d["added"]:   print(f"    + {a!r}")
        print("  → Edit make_slides.py manually for these.")

    return applied


# ── AI instruction reporter ───────────────────────────────────────────────────

def report_ai_notes(slides: list[dict]) -> list[tuple[int, str, list[str]]]:
    """Return (slide_num, title, [instructions]) for slides with AI: notes."""
    return [
        (s["slide_num"], s["title"], s["ai_notes"])
        for s in slides
        if s["ai_notes"]
    ]


def strip_ai_notes_from_script(ai_items: list[tuple[int, str, list[str]]]) -> None:
    """Remove AI: lines from notes= strings in make_slides.py."""
    source = SCRIPT.read_text()
    for _slide_num, _title, instructions in ai_items:
        for instr in instructions:
            # Match   AI: <instruction>  with any surrounding whitespace / newline escapes
            pattern = r'\\n\s*AI:\s*' + re.escape(instr) + r'\s*'
            source, n = re.subn(pattern, '', source)
            if n == 0:
                # Try without leading \n
                pattern2 = r'AI:\s*' + re.escape(instr) + r'\s*\\n?'
                source, _ = re.subn(pattern2, '', source)
    SCRIPT.write_text(source)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description=__doc__,
                                     formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("--apply", action="store_true",
                        help="Apply patches to make_slides.py (default: dry-run)")
    args = parser.parse_args()

    if not PPTX.exists():
        sys.exit(f"ERROR: {PPTX} not found — run from the project root")
    if not SCRIPT.exists():
        sys.exit(f"ERROR: {SCRIPT} not found")

    # ── Step 1: save user's edited PPTX ──────────────────────────────────────
    shutil.copy(PPTX, EDITED)
    print(f"Saved edited PPTX → {EDITED}")

    # ── Step 2: regenerate baseline ───────────────────────────────────────────
    print("Regenerating baseline from make_slides.py …")
    result = subprocess.run(
        [sys.executable, str(SCRIPT)],
        capture_output=True, text=True
    )
    if result.returncode != 0:
        sys.exit(f"ERROR: make_slides.py failed:\n{result.stderr}")
    shutil.copy(PPTX, BASELINE)
    print(f"Baseline saved → {BASELINE}")

    # ── Step 3: extract and diff ──────────────────────────────────────────────
    baseline_slides = extract_slides(BASELINE)
    edited_slides   = extract_slides(EDITED)

    text_diffs = find_text_diffs(baseline_slides, edited_slides)
    ai_items   = report_ai_notes(edited_slides)

    # ── Report / apply text diffs ─────────────────────────────────────────────
    if text_diffs:
        print(f"\n{'─'*60}")
        print(f"TEXT CHANGES ({len(text_diffs)} slide(s)):")
        for d in text_diffs:
            print(f"\n  Slide {d['slide_num']} — {d['title']!r}")
            for r in d["removed"]: print(f"    - {r!r}")
            for a in d["added"]:   print(f"    + {a!r}")
        if args.apply:
            applied = patch_script(text_diffs)
            print(f"\nPatched make_slides.py:")
            for a in applied:
                print(f"  ✓ {a}")
    else:
        print("\nNo visible text changes.")

    # ── Report AI instructions ─────────────────────────────────────────────────
    if ai_items:
        print(f"\n{'─'*60}")
        print(f"AI INSTRUCTIONS ({sum(len(i) for _,_,i in ai_items)} item(s)):")
        for slide_num, title, instructions in ai_items:
            print(f"\n  Slide {slide_num} — {title!r}")
            for instr in instructions:
                print(f"    AI: {instr}")
        if not args.apply:
            print("\n  (dry-run: pass --apply to act on these)")
        else:
            print("\n  → Claude will action these instructions.")
            print("    AI: lines will be stripped from make_slides.py after Claude acts on them.")
    else:
        print("\nNo AI instructions found.")

    # ── After --apply: regenerate PPTX to keep it in sync with make_slides.py ──
    if args.apply and (text_diffs or ai_items):
        print(f"\n{'─'*60}")
        print("Regenerating PPTX from updated make_slides.py …")
        result2 = subprocess.run(
            [sys.executable, str(SCRIPT)],
            capture_output=True, text=True
        )
        if result2.returncode != 0:
            print(f"WARNING: make_slides.py failed after patch:\n{result2.stderr}")
        else:
            print(f"PPTX regenerated — {PPTX} is now in sync with make_slides.py.")

    if not args.apply:
        print(f"\n{'─'*60}")
        print("Dry-run complete.  Re-run with --apply to patch make_slides.py.")


if __name__ == "__main__":
    main()
