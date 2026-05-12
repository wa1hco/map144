#!/home/jeff/ham/map144/.venv/bin/python3
"""Generate map144 presentation as map144_presentation.pptx.

Run from the project root:
    pip install python-pptx        # one-time; not a runtime dep
    python docs/make_slides.py

Produces docs/map144_presentation.pptx  (open in LibreOffice Impress or PowerPoint).

Slides use proper Title and Content placeholder shapes so Tab-to-indent works
in LibreOffice. Screenshot placeholders are light blue boxes — replace with
actual screenshots before presenting.

Color scheme: white background, navy titles, blue accent. Light and readable
in bright rooms; prints cleanly.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE
from pathlib import Path
from PIL import Image as _PILImage
import subprocess
from lxml import etree as _etree
from pptx.oxml.ns import qn as _qn

# ── color palette (light theme) ─────────────────────────────────────────────
BG           = RGBColor(0xFF, 0xFF, 0xFF)   # white background
NAVY         = RGBColor(0x1B, 0x3A, 0x6B)   # dark navy — titles
ACCENT       = RGBColor(0x27, 0x6F, 0xBF)   # medium blue — accents, section headers
BODY         = RGBColor(0x1A, 0x1A, 0x2E)   # near-black — body text level 0
MUTED        = RGBColor(0x44, 0x44, 0x55)   # dark grey — sub-bullet text
LGREY        = RGBColor(0x88, 0x88, 0x99)   # light grey — captions, notes text
SECTION_BG   = RGBColor(0x1B, 0x3A, 0x6B)   # navy — section break slide background
SECTION_FG   = RGBColor(0xFF, 0xFF, 0xFF)   # white — section break text
PLACEHOLDER  = RGBColor(0xD6, 0xEA, 0xF8)   # pale blue — screenshot placeholder fill
PLACEHOLDER_BORDER = RGBColor(0x27, 0x6F, 0xBF)
MONO_BG      = RGBColor(0xF0, 0xF4, 0xF8)   # very light blue — code block background
MONO_FG      = RGBColor(0x1B, 0x3A, 0x6B)   # navy — code text

TITLE_FONT = "Calibri"
BODY_FONT  = "Calibri"
MONO_FONT  = "Courier New"

W = Inches(13.333)
H = Inches(7.5)

_DOCS_DIR    = Path(__file__).parent
_FIGURES_DIR = _DOCS_DIR / "figures"          # all PNG assets live here
_FIGURES_DIR.mkdir(exist_ok=True)

_PIPELINE_PNG = _FIGURES_DIR / "map144_pipeline.png"

_PIPELINE_DOT = """\
digraph pipeline {
    rankdir=LR
    graph [bgcolor="#FFFFFF", pad=0.45, nodesep=0.30, ranksep=0.50]
    node  [shape=box, style="filled,rounded", fillcolor="#D6EAF8",
           color="#276FBF", fontname="Helvetica", fontsize=12,
           margin="0.15,0.08"]
    edge  [color="#1B3A6B", arrowsize=0.8, penwidth=1.5, fontname="Helvetica",
           fontsize=9]

    // ── Processing nodes ──────────────────────────────────────────────────
    src  [label="IQ Source\\n48 kHz"]
    nb   [label="Noise\\nBlanker"]
    ch   [label="Polyphase\\nChannelizer\\n48 × 1 kHz"]
    sq   [label="Squarer\\n(x\\u00B2)"]
    det  [label="Metric\\nDetector"]
    gate [label="Coincidence\\nGate"]
    bpf  [label="Decimate+Mix\\n+BPF\\n48→12 kHz IQ"]
    spd  [label="SPD\\nDecoder", fillcolor="#D5F5E3"]
    jt9  [label="jt9 Decoder\\n(fallback)", fillcolor="#D5F5E3"]
    dl   [label="Decode List", fillcolor="#D5F5E3"]

    // ── Storage ───────────────────────────────────────────────────────────
    ring [label="IQ Ring Buffer\\n30 s", shape=cylinder,
          fillcolor="#FEF9E7", color="#8E6914"]

    // ── Display windows & outputs (row 2) ─────────────────────────────────
    iqwin [label="IQ / Noise\\nBlanker Window", shape=box,
           fillcolor="#FCF3CF", color="#8E6914"]
    fg    [label="Fast Graph\\nWaterfall",      shape=box,
           fillcolor="#FCF3CF", color="#8E6914"]
    dp    [label="Tone Detection\\nSNR Plot",   shape=box,
           fillcolor="#FCF3CF", color="#8E6914"]
    rep   [label="Reporting",                   shape=box,
           fillcolor="#FADBD8", color="#8E6914"]

    // ── Main pipeline ─────────────────────────────────────────────────────
    src -> nb -> ch -> sq -> det -> gate -> bpf -> spd
    spd -> dl [label=" decoded"]
    spd -> jt9 [style=dashed, label=" SPD miss\\n(fallback)"]
    jt9 -> dl

    // ── IQ ring buffer (for jt9 fallback) ─────────────────────────────────
    ch   -> ring [style=dashed, constraint=false, label=" write"]
    ring -> jt9  [constraint=false, label=" 12 kHz\\naudio"]

    // ── IQ/Noise Blanker Window: before & after blanker ───────────────────
    src  -> iqwin [style=dashed, color="#888899", constraint=false]
    nb   -> iqwin [style=dashed, color="#888899", constraint=false]

    // ── Fast Graph: waterfall from channelizer only ───────────────────────
    ch   -> fg [style=dashed, color="#888899", constraint=false]

    // ── Tone Detection SNR Plot ───────────────────────────────────────────
    sq   -> dp [style=dashed, color="#888899", constraint=false,
                label=" background"]
    gate -> dp [style=dashed, color="#888899", constraint=false,
                label=" circles"]
    spd  -> dp [style=dashed, color="#888899", constraint=false,
                label=" color\\ncodes"]
    jt9  -> dp [style=dashed, color="#888899", constraint=false,
                label=" color\\ncodes"]

    // ── Reporting: decoded messages forwarded externally ──────────────────
    spd  -> rep [style=dashed, color="#888899", constraint=false]
    jt9  -> rep [style=dashed, color="#888899", constraint=false]

    // ── Row-2 placement: anchor each display node below its source ─────────
    { rank=same; nb;   iqwin }
    { rank=same; ch;   fg    }
    { rank=same; sq;   ring  }
    { rank=same; gate; dp    }
    { rank=same; dl;   rep   }
}
"""

_SPD_DETAIL_PNG = _FIGURES_DIR / "map144_spd_detail.png"

_SPD_DETAIL_DOT = """\
digraph spd_detail {
    rankdir=LR
    graph [bgcolor="#FFFFFF", pad=0.45, nodesep=0.35, ranksep=0.55]
    node  [shape=box, style="filled,rounded", fillcolor="#D6EAF8",
           color="#276FBF", fontname="Helvetica", fontsize=12,
           margin="0.18,0.10"]
    edge  [color="#1B3A6B", arrowsize=0.8, penwidth=1.5, fontname="Helvetica",
           fontsize=9]

    iq   [label="Complex IQ\n12 kHz\n(channelizer)"]
    bpf  [label="Audio BPF\n\xb11200 Hz\nfiltfilt 129-tap"]
    sq   [label="Squared-spectrum\nsliding window\n72 ms hop"]
    det  [label="Candidate\ndetection\n(metric threshold)"]
    fc   [label="Freq search\n\xb1200 Hz grid"]
    sync [label="Sync correlation\n(threshold 1.3)"]
    avg  [label="Frame averaging\n(navg frames)"]
    bp   [label="BP decoder\nbpdecode128_90"]
    ldpc [label="LDPC check\n+ CRC-12"]
    out  [label="Decoded\nmessage", fillcolor="#D5F5E3"]
    miss [label="SPD miss\n\u2192 jt9 fallback", shape=box,
          style="filled,rounded", fillcolor="#FADBD8", color="#8E6914"]

    iq -> bpf -> sq -> det -> fc -> sync -> avg -> bp -> ldpc -> out
    det  -> miss [style=dashed, color="#CC4444", label=" no candidate"]
    sync -> miss [style=dashed, color="#CC4444", label=" sync fail"]
    ldpc -> miss [style=dashed, color="#CC4444", label=" CRC fail"]
}
"""

def _make_pipeline_png():
    """Render _PIPELINE_DOT to _PIPELINE_PNG via the graphviz dot CLI."""
    dot_path = _DOCS_DIR / "map144_pipeline.dot"
    dot_path.write_text(_PIPELINE_DOT)
    result = subprocess.run(
        ["dot", "-Tpng", "-Gdpi=150", str(dot_path), "-o", str(_PIPELINE_PNG)],
        capture_output=True, text=True,
    )
    if result.returncode != 0:
        print(f"[warning] dot failed: {result.stderr.strip()}")
    else:
        print(f"Generated: {_PIPELINE_PNG}")


def _make_spd_detail_png():
    """Render _SPD_DETAIL_DOT to _SPD_DETAIL_PNG via the graphviz dot CLI."""
    dot_path = _DOCS_DIR / "map144_spd_detail.dot"
    dot_path.write_text(_SPD_DETAIL_DOT)
    result = subprocess.run(
        ["dot", "-Tpng", "-Gdpi=150", str(dot_path), "-o", str(_SPD_DETAIL_PNG)],
        capture_output=True, text=True,
    )
    if result.returncode != 0:
        print(f"[warning] dot failed: {result.stderr.strip()}")
    else:
        print(f"Generated: {_SPD_DETAIL_PNG}")


# ── Helpers ───────────────────────────────────────────────────────────────────

def _set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_run(run, color, size_pt, bold=False, italic=False, font=BODY_FONT):
    run.font.name      = font
    run.font.size      = Pt(size_pt)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_textbox(slide, text, left, top, width, height,
                color=BODY, size=20, bold=False, italic=False,
                align=PP_ALIGN.LEFT, font=BODY_FONT):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text = text
    _set_run(r, color, size, bold=bold, italic=italic, font=font)
    return txb


def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_pt=0):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width     = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def add_image_placeholder(slide, label, left, top, width, height):
    shape = add_rect(slide, left, top, width, height,
                     PLACEHOLDER, PLACEHOLDER_BORDER, line_pt=1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r  = p.add_run()
    r.text = f"[ Screenshot: {label} ]"
    _set_run(r, ACCENT, 16, italic=True)
    return shape


def _add_data_table(slide, headers, rows, col_widths_in, left, top,
                    row_h_in=0.40, base_font_pt=14):
    """Add a navy-header striped data table to a slide.

    headers       – list of column header strings
    rows          – list of lists (one list per data row)
    col_widths_in – list of column widths in inches (must match len(headers))
    left, top     – position in EMU
    """
    n_rows = len(rows) + 1   # +1 for header
    n_cols = len(headers)
    tbl = slide.shapes.add_table(
        n_rows, n_cols,
        int(left), int(top),
        int(Inches(sum(col_widths_in))),
        int(Inches(row_h_in * n_rows)),
    ).table
    for ci, w in enumerate(col_widths_in):
        tbl.columns[ci].width = int(Inches(w))
    # Header row
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0] if p.runs else p.add_run()
        r.font.name  = BODY_FONT
        r.font.size  = Pt(base_font_pt)
        r.font.bold  = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    # Data rows
    for ri, row_data in enumerate(rows):
        for ci, cell_text in enumerate(row_data):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(cell_text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                RGBColor(0xF0, 0xF4, 0xF8) if ri % 2 == 0
                else RGBColor(0xFF, 0xFF, 0xFF)
            )
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.runs[0] if p.runs else p.add_run()
            r.font.name  = BODY_FONT
            r.font.size  = Pt(base_font_pt)
            r.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)


# ── Slide layouts in the default python-pptx template ────────────────────────
# 0  Title Slide          — title + subtitle placeholders
# 1  Title and Content    — title + body (bullet)
# 2  Title and Two Content
# 3  Title Only
# 6  Blank

def _blank(prs, bg=BG):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, bg)
    return slide


def _title_content(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_bg(slide)
    return slide


def _two_content(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[3])  # Title Only — we add two boxes manually
    _set_bg(slide)
    return slide


def _set_title_placeholder(slide, text, color=NAVY, size=32):
    title_ph = slide.shapes.title
    title_ph.text = text
    for para in title_ph.text_frame.paragraphs:
        for run in para.runs:
            _set_run(run, color, size, bold=True)


def _fill_body_placeholder(slide, items, ph_idx=1, base_size=20, mono_items=None):
    """Fill a body placeholder with (level, text) items.

    mono_items: set of item indices to render in monospace font.
    level 0 → bold body color; level 1+ → muted, slightly smaller.
    """
    mono_items = mono_items or set()
    ph = slide.placeholders[ph_idx]
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    first = True
    for idx, (level, text) in enumerate(items):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        is_mono = idx in mono_items
        r = p.add_run()
        r.text = text
        color = BODY if level == 0 else MUTED
        size  = base_size - level * 2
        _set_run(r, color, size, bold=(level == 0 and not is_mono),
                 font=MONO_FONT if is_mono else BODY_FONT)


def _fit_placeholders(slide):
    """Resize title and body placeholders to full slide width.

    Called on every _title_content slide so bullets use all available horizontal
    space.  bullets_image_slide uses _blank + manual textboxes and is unaffected.
    """
    _margin = Inches(0.4)
    _full_w = W - 2 * _margin
    _footer_top = H - _FOOTER_H

    title_ph = slide.shapes.title
    title_ph.left   = int(_margin)
    title_ph.top    = int(Inches(0.15))
    title_ph.width  = int(_full_w)
    title_ph.height = int(Inches(0.85))

    body_ph = slide.placeholders[1]
    body_ph.left   = int(_margin)
    body_ph.top    = int(Inches(1.1))
    body_ph.width  = int(_full_w)
    body_ph.height = int(_footer_top - Inches(1.25))


# ── Slide builders ────────────────────────────────────────────────────────────

def title_slide(prs, title, subtitle, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _set_bg(slide)

    # Accent bar at top
    add_rect(slide, 0, Inches(0.0), W, Pt(10), ACCENT)

    title_ph    = slide.shapes.title
    subtitle_ph = slide.placeholders[1]

    title_ph.text = title
    for para in title_ph.text_frame.paragraphs:
        for run in para.runs:
            _set_run(run, NAVY, 48, bold=True)

    subtitle_ph.text = subtitle
    for para in subtitle_ph.text_frame.paragraphs:
        for run in para.runs:
            _set_run(run, ACCENT, 28)

    # Author strip at bottom
    strip = add_rect(slide, 0, Inches(6.8), W, Inches(0.7), NAVY)
    add_textbox(slide, "Jeff Millar  WA1HCO",
                Inches(0.5), Inches(6.82), W - Inches(1), Inches(0.5),
                color=SECTION_FG, size=18, align=PP_ALIGN.CENTER)

    if notes:
        add_notes(slide, notes)
    return slide


def section_slide(prs, section_num, title, subtitle="", audience="", notes=""):
    """Navy background section-break slide — stands out from content slides."""
    slide = _blank(prs, bg=SECTION_BG)

    # Left accent stripe in lighter blue
    add_rect(slide, 0, 0, Inches(0.3), H, ACCENT)

    add_textbox(slide, f"Section {section_num}",
                Inches(0.6), Inches(1.8), Inches(12), Inches(0.7),
                color=ACCENT, size=22, bold=True)
    add_textbox(slide, title,
                Inches(0.6), Inches(2.5), Inches(12), Inches(1.6),
                color=SECTION_FG, size=48, bold=True)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.6), Inches(4.2), Inches(12), Inches(0.8),
                    color=RGBColor(0xCC, 0xDD, 0xEE), size=24, italic=True)
    if audience:
        add_textbox(slide, f"Audience: {audience}",
                    Inches(0.6), Inches(6.4), Inches(12), Inches(0.6),
                    color=ACCENT, size=16)

    note_text = (f"[Audience: {audience}]\n\n" if audience else "") + (notes or "")
    if note_text:
        add_notes(slide, note_text)
    return slide


def content_slide(prs, title, items, notes="", mono_items=None):
    """Standard title + bullet body slide using real Content placeholder."""
    slide = _title_content(prs)
    add_rect(slide, 0, 0, W, Pt(8), ACCENT)          # thin top bar
    _fit_placeholders(slide)
    _set_title_placeholder(slide, title)
    _fill_body_placeholder(slide, items, mono_items=mono_items)
    if notes:
        add_notes(slide, notes)
    return slide


def _add_picture_fit(slide, png_path, left, top, max_w, max_h, align='center'):
    """Insert a picture within the bounding box, preserving aspect ratio.

    align: 'center' (default), 'right' (flush to right edge of box),
           or 'left' (flush to left edge of box).
    """
    with _PILImage.open(png_path) as im:
        iw, ih = im.size
    scale = min(max_w / iw, max_h / ih)
    pw = int(iw * scale)
    ph = int(ih * scale)
    if align == 'right':
        cx = left + (max_w - pw)
    elif align == 'left':
        cx = left
    else:
        cx = left + (max_w - pw) // 2
    cy = top + (max_h - ph) // 2
    slide.shapes.add_picture(str(png_path), cx, cy, pw, ph)


def _apply_textbox_bullet(p, level):
    """Add a bullet character and hanging indent to a textbox paragraph.

    Plain textboxes don't inherit theme bullet styles, so we write the XML
    directly.  Level 0 gets •, level 1 gets –, level 2+ gets ◦.
    """
    _HANG    = 342900   # 3/8" in EMU — hanging indent overhang
    _CHARS   = {0: '\u2022', 1: '\u2013', 2: '\u25e6'}
    pPr = p._p.get_or_add_pPr()
    pPr.set('marL',   str(_HANG + level * _HANG * 2))
    pPr.set('indent', str(-_HANG))
    for tag in ('a:buNone', 'a:buChar', 'a:buAutoNum'):
        el = pPr.find(_qn(tag))
        if el is not None:
            pPr.remove(el)
    buChar = _etree.SubElement(pPr, _qn('a:buChar'))
    buChar.set('char', _CHARS.get(level, '\u2022'))


def image_slide(prs, title, label, caption="", notes="", png=None):
    """Full-width image (or placeholder) slide — no bullet text."""
    slide = _blank(prs)
    add_rect(slide, 0, 0, W, Pt(8), ACCENT)
    add_textbox(slide, title,
                Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.85),
                color=NAVY, size=32, bold=True)
    img_h = Inches(5.4) if caption else Inches(6.1)
    png_path = _FIGURES_DIR / png if png else None
    if png_path and png_path.exists():
        _add_picture_fit(slide, png_path,
                         int(Inches(0.4)), int(Inches(1.1)),
                         int(Inches(12.5)), int(img_h))
    else:
        add_image_placeholder(slide, label,
                              Inches(0.4), Inches(1.1), Inches(12.5), img_h)
    if caption:
        add_textbox(slide, caption,
                    Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.55),
                    color=LGREY, size=15, italic=True)
    if notes:
        add_notes(slide, notes)
    return slide


def bullets_image_slide(prs, title, items, label, caption="", notes="", png=None,
                        img_right_margin=Inches(0.4), img_align='center',
                        img_top=None, img_h=None,
                        img_left=None, txt_w=None):
    """Left column: bullet text.  Right column: screenshot or placeholder.

    img_right_margin: gap between image right edge and slide edge (default 0.4").
    img_align:        'center' | 'right' — horizontal alignment within the box.
    img_top:          override image top position in EMU (default: Inches(1.05)).
    img_h:            override image box height in EMU (default: fills to footer).
    img_left:         override image column left edge in EMU (default: Inches(6.1)).
    txt_w:            override text column width in EMU (default: Inches(5.5)).
    Image is inserted before text shapes so text is always on top.
    """
    slide = _blank(prs)
    add_rect(slide, 0, 0, W, Pt(8), ACCENT)

    # ── Right: image (added first → behind text) ──────────────────────────
    _IMG_LEFT  = img_left if img_left is not None else Inches(6.1)
    _IMG_W     = W - _IMG_LEFT - img_right_margin
    _txt_top   = Inches(1.05)
    _def_img_h = Inches(5.5) if not caption else Inches(5.0)
    _i_top     = img_top if img_top is not None else int(_txt_top)
    _i_h       = img_h   if img_h   is not None else int(_def_img_h)
    png_path = _FIGURES_DIR / png if png else None
    if png_path and png_path.exists():
        _add_picture_fit(slide, png_path,
                         int(_IMG_LEFT), _i_top, int(_IMG_W), _i_h,
                         align=img_align)
    else:
        add_image_placeholder(slide, label, _IMG_LEFT, _i_top, _IMG_W, _i_h)

    # ── Title and bullets (added after image → on top) ────────────────────
    add_textbox(slide, title,
                Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.75),
                color=NAVY, size=28, bold=True)

    _LEFT_W = txt_w if txt_w is not None else Inches(5.5)
    _BODY_H = Inches(5.5) if not caption else Inches(5.0)
    txb = slide.shapes.add_textbox(Inches(0.4), _txt_top, _LEFT_W, _BODY_H)
    tf  = txb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    first = True
    for level, text in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = text
        _set_run(r, BODY if level == 0 else MUTED,
                 18 - level * 2, bold=(level == 0))
        _apply_textbox_bullet(p, level)

    if caption:
        add_textbox(slide, caption,
                    Inches(0.4), Inches(6.65), Inches(12.5), Inches(0.6),
                    color=LGREY, size=14, italic=True)
    if notes:
        add_notes(slide, notes)
    return slide


def stacked_image_slide(prs, title, items, label, notes="", png=None):
    """Bullets full-width on top; image at the bottom just above the footer.

    Image is scaled to fit the full content width (margin on both sides and
    right edge).  Its bottom sits a small gap above the footer strip.  The
    text box is sized to exactly fill the space from the title to the image
    top — no fixed split ratio, no wasted whitespace.

    Bullet style matches content_slide (_fill_body_placeholder: 20 pt base,
    BODY/MUTED colours, bold level-0).
    """
    slide = _blank(prs)
    add_rect(slide, 0, 0, W, Pt(8), ACCENT)

    _margin     = Inches(0.4)
    _full_w     = W - 2 * _margin
    _txt_top    = Inches(1.05)
    _gap        = Inches(0.12)                # gap between text box and image
    _img_bottom = _FOOTER_TOP - Inches(0.12)  # gap above footer strip
    _min_txt_h  = Inches(1.0)                 # minimum height reserved for bullets

    # ── Compute image geometry from actual pixel dimensions ───────────────
    png_path = _FIGURES_DIR / png if png else None
    if png_path and png_path.exists():
        with _PILImage.open(png_path) as im:
            iw, ih = im.size
        _max_img_h = int(_img_bottom - _txt_top - _min_txt_h - _gap)
        scale     = min(_full_w / iw, _max_img_h / ih)
        pw        = int(iw * scale)
        ph        = int(ih * scale)
        _img_top  = int(_img_bottom) - ph
        _img_left = int(_margin) + (int(_full_w) - pw) // 2
    else:
        ph        = int(Inches(3.5))
        pw        = int(_full_w)
        _img_top  = int(_img_bottom) - ph
        _img_left = int(_margin)

    # ── Image (added first → behind text) ────────────────────────────────
    if png_path and png_path.exists():
        slide.shapes.add_picture(str(png_path), _img_left, _img_top, pw, ph)
    else:
        add_image_placeholder(slide, label,
                              int(_margin), _img_top,
                              int(_full_w), int(_img_bottom) - _img_top)

    # ── Title and bullets (added after image → on top) ────────────────────
    add_textbox(slide, title,
                _margin, Inches(0.15), _full_w, Inches(0.75),
                color=NAVY, size=28, bold=True)

    _txt_h = _img_top - int(_txt_top) - int(_gap)
    txb = slide.shapes.add_textbox(int(_margin), int(_txt_top),
                                   int(_full_w), max(_txt_h, int(_min_txt_h)))
    tf  = txb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    first = True
    for level, text in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = text
        _set_run(r, BODY if level == 0 else MUTED,
                 20 - level * 2, bold=(level == 0))
        _apply_textbox_bullet(p, level)

    if notes:
        add_notes(slide, notes)
    return slide


def two_col_slide(prs, title, left_items, right_items,
                  left_title="", right_title="", notes=""):
    slide = _blank(prs)
    add_rect(slide, 0, 0, W, Pt(8), ACCENT)
    add_textbox(slide, title,
                Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.85),
                color=NAVY, size=32, bold=True)

    def _col_box(slide, items, left, col_title):
        top   = Inches(1.1)
        width = Inches(5.9)
        if col_title:
            add_textbox(slide, col_title, left, top, width, Inches(0.55),
                        color=ACCENT, size=22, bold=True)
            top += Inches(0.6)
        txb = slide.shapes.add_textbox(left, top, width, Inches(5.6))
        tf  = txb.text_frame
        tf.word_wrap = True
        first = True
        for level, text in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.level = level
            r = p.add_run()
            r.text = text
            _set_run(r, BODY if level == 0 else MUTED,
                     20 - level * 2, bold=(level == 0))

    _col_box(slide, left_items,  Inches(0.4),  left_title)
    _col_box(slide, right_items, Inches(6.9), right_title)

    if notes:
        add_notes(slide, notes)
    return slide


def quote_slide(prs, quote, attribution="", notes=""):
    slide = _blank(prs)
    add_rect(slide, 0, 0, W, Pt(8), ACCENT)
    add_textbox(slide, "\u201c",
                Inches(0.3), Inches(0.6), Inches(2), Inches(2),
                color=ACCENT, size=120, bold=True)
    add_textbox(slide, quote,
                Inches(1.2), Inches(1.4), Inches(11.2), Inches(4.2),
                color=NAVY, size=28, italic=True)
    if attribution:
        add_textbox(slide, f"\u2014 {attribution}",
                    Inches(1.2), Inches(5.7), Inches(11.2), Inches(0.8),
                    color=ACCENT, size=20)
    if notes:
        add_notes(slide, notes)
    return slide


def code_slide(prs, title, code_text, caption="", notes=""):
    slide = _blank(prs)
    add_rect(slide, 0, 0, W, Pt(8), ACCENT)
    add_textbox(slide, title,
                Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.85),
                color=NAVY, size=32, bold=True)
    box_h = Inches(5.4) if caption else Inches(6.1)
    shape = add_rect(slide, Inches(0.4), Inches(1.1), Inches(12.5), box_h,
                     MONO_BG, LGREY, line_pt=1.0)
    tf = shape.text_frame
    tf.word_wrap = False
    first = True
    for line in code_text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = line
        _set_run(r, MONO_FG, 16, font=MONO_FONT)
    if caption:
        add_textbox(slide, caption,
                    Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.55),
                    color=LGREY, size=15, italic=True)
    if notes:
        add_notes(slide, notes)
    return slide


# ── Footer ────────────────────────────────────────────────────────────────────

_FOOTER_H    = Inches(0.28)
_FOOTER_TOP  = H - _FOOTER_H
_FOOTER_DATE = "April 2026"

def _add_footer(slide, page_num):
    """Thin navy strip at slide bottom: date left, page-number right."""
    add_rect(slide, 0, _FOOTER_TOP, W, _FOOTER_H, NAVY)
    add_textbox(slide, _FOOTER_DATE,
                Inches(0.3), _FOOTER_TOP, Inches(4), _FOOTER_H,
                color=RGBColor(0xFF, 0xFF, 0xFF), size=11, align=PP_ALIGN.LEFT)
    add_textbox(slide, str(page_num),
                Inches(0.0), _FOOTER_TOP, W - Inches(0.3), _FOOTER_H,
                color=RGBColor(0xFF, 0xFF, 0xFF), size=11, align=PP_ALIGN.RIGHT)


# ── Build the deck ────────────────────────────────────────────────────────────

def build():
    _make_pipeline_png()
    _make_spd_detail_png()

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    # ── Title ─────────────────────────────────────────────────────────────────
    title_slide(prs,
        "MAP144",
        "Wideband MSK144 Meteor Scatter Monitor",
        notes="Introduction slide. Adjust callsign / date before presenting.")

    # ── Agenda ────────────────────────────────────────────────────────────────
    content_slide(prs, "Agenda", [
        (0, "1 · Why MAP144?"),
        (0, "2 · Background — meteor scatter and MSK144"),
        (0, "3 · Features of MAP144"),
        (0, "4 · Files generated during operation"),
        (0, "5 · Installation"),
        (0, "6 · How to operate MAP144"),
        (0, "7 · System architecture (technical deep-dive)"),
        (0, "8 · Building the project with AI assistance"),
        (0, "9 · Performance engineering — iterative profiling and optimization"),
        (0, "10 · Decoder sensitivity investigation"),
        (0, "11 · MAP144 vs WSJT-X real-world comparison"),
        (0, "12 · Doppler measurements & design opportunities"),
        (0, "13 · Support and resources"),
    ], notes="Overview of the talk. Skip or reorder sections to suit the audience. "
             "Section 12 is the original-research half — skip for general audiences.")

    content_slide(prs, "Why MAP144?", [
        (0, "Normal MSK144 activity is concentrated on the calling frequency"),
        (1, "Little interference — few stations active at any one time"),
        (1, "WSJT-X handles this well"),
        (0, "During contests, increased activity"),
        (1, "Operators spread out from the calling frequency"),
        (1, "Scheduled contacts avoid the calling frequency"),
        (1, "Activity is happening across many frequencies simultaneously"),
        (1, "MAP144 Decodes all MSK144 activity across ±17 kHz from the calling frequency"),
        (1, "Contest operator sees who is active, on what frequency, at what SNR"),
        (1, "Decoded contacts get posted to GridTracker, N1MM, PSKreporter"),
        (1, "Operators act on that information — schedule a contact or tail end"),
        (0, "Inspired by MAP65 (EME) and QMAP (Q65) — same philosophy, new mode"),
    ], notes="This is the original design motivation — contest situational awareness. "
             "The contest operator runs WSJT-X for their own contacts and MAP144 as "
             "a passive band monitor feeding GridTracker.")

    # ═══════════════════════════════════════════════════════════════════════════
    # SECTION 2 — BACKGROUND
    # ═══════════════════════════════════════════════════════════════════════════
    section_slide(prs, 2, "Background",
        subtitle="Meteor scatter and the MSK144 protocol",
        audience="Everyone",
        notes="Good starting point for a general ham radio audience.")

    # ── What is Meteor Scatter — bullets + shower table ───────────────────────
    _ms_slide = _title_content(prs)
    add_rect(_ms_slide, 0, 0, W, Pt(8), ACCENT)
    _fit_placeholders(_ms_slide)
    _set_title_placeholder(_ms_slide, "What is Meteor Scatter?")
    _fill_body_placeholder(_ms_slide, [
        (0, "Meteors burn up in the ionosphere at ~80–120 km altitude"),
        (1, "Ionised trail reflects VHF signals — briefly, unpredictably"),
        (0, "Useful for 50, 144, and 220 MHz communication beyond line-of-sight"),
        (1, "Typical path: 500–2000 km with no repeater or satellite"),
        (0, "Pings last milliseconds to a few seconds"),
        (1, "Random arrival — you must monitor continuously"),
        (0, "Top annual showers:"),
    ])
    # Table: Shower | Peak | ZHR
    _ROWS = [
        ("Shower",       "Peak date",   "ZHR"),
        ("Geminids",     "Dec 13–14",   "~150"),
        ("Perseids",     "Aug 11–13",   "~100"),
        ("Quadrantids",  "Jan 3–4",     "~120 (narrow, ~6 hr peak)"),
        ("Leonids",      "Nov 17–18",   "~15 typical; storms ~1000+"),
        ("Eta Aquariids","May 5–6",     "~50–70"),
    ]
    from pptx.util import Inches as _In
    _tbl = _ms_slide.shapes.add_table(
        len(_ROWS), 3,
        _In(1.0), _In(4.9), _In(11.3), _In(1.85)
    ).table
    _tbl.columns[0].width = _In(2.4)
    _tbl.columns[1].width = _In(2.0)
    _tbl.columns[2].width = _In(6.9)
    from pptx.util import Pt as _Pt
    from pptx.dml.color import RGBColor as _RGB
    for r_i, row_data in enumerate(_ROWS):
        for c_i, cell_text in enumerate(row_data):
            cell = _tbl.cell(r_i, c_i)
            cell.text = cell_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                _RGB(0x1B, 0x3A, 0x6B) if r_i == 0   # header: navy
                else (_RGB(0xF0, 0xF4, 0xF8) if r_i % 2 == 0 else _RGB(0xFF, 0xFF, 0xFF))
            )
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.font.name  = BODY_FONT
            run.font.size  = _Pt(14)
            run.font.bold  = (r_i == 0)
            run.font.color.rgb = _RGB(0xFF, 0xFF, 0xFF) if r_i == 0 else _RGB(0x1A, 0x1A, 0x2E)
    add_notes(_ms_slide, "Set the stage for why automated monitoring is valuable. "
              "Quadrantids have a very narrow peak — often missed by casual operators. "
              "Leonids are unpredictable but historic storms produce extraordinary 2m activity.")

    content_slide(prs, "The Propagation Path", [
        (0, "Meteor trail forms at ~90 km altitude"),
        (1, "Typical path length: 500–2000 km  —  trail midpoint roughly equidistant"),
        (0, "Trail density — linear electron density (electrons per meter of trail)"),
        (1, "Volume density is not useful — trail diameter expands over time"),
        (1, "Underdense (<10¹⁴ e/m): trail transparent to the wave"),
        (2, "Each electron scatters independently — signal adds in power"),
        (1, "Overdense (>10¹⁴ e/m): trail acts as a reflecting cylinder"),
        (2, "Signal reflects off the surface — stronger but shorter-lived"),
    ], notes="Keep this slide conceptual. The underdense/overdense distinction "
             "explains why pings vary so much in character and duration.")

    content_slide(prs, "Meteor Scatter vs. Operating Frequency", [
        (0, "Most key quantities scale as 1/f²  —  higher frequency = harder"),
        (0, "Received signal power (underdense trails)"),
        (1, "Power ∝ 1/f²  —  doubling frequency quarters the received power"),
        (0, "Overdense threshold"),
        (1, "Required electron line density ∝ f²  —  rarer meteors qualify at higher freq"),
        (0, "Trail duration"),
        (1, "Decay time constant ∝ 1/f²  —  trails fade faster at higher frequencies"),
        (1, "Diffusion spreads the trail; at higher f the density drops below threshold sooner"),
        (0, "Faraday rotation"),
        (1, "Rotation angle ∝ 1/f²  —  much less significant at 144 MHz than 50 MHz"),
        (1, "At 432 MHz polarization rotation is nearly negligible"),
        (0, "Practical bands: 50 MHz (best), 144 MHz (common), 222 MHz (marginal), 432 MHz (rare)"),
    ], notes="The 1/f² scaling comes from the plasma physics of free electron scattering. "
             "Good slide for a technical audience — unifies several phenomena with one rule. "
             "The Faraday rotation point motivates dual-pol being most valuable at 6 m.")

    content_slide(prs, "Underdense vs. Overdense — Why It Matters", [
        (0, "Underdense trail — the common case"),
        (1, "Short ping, fades in 50–300 ms"),
        (1, "MSK144 may capture only one or two 72 ms frames"),
        (1, "Decode success depends heavily on SNR in those few frames"),
        (1, "Missed decodes are common — signal was real but too brief"),
        (0, "Overdense trail — the dramatic case"),
        (1, "Longer, initially stronger signal — visible on the waterfall"),
        (1, "Trail distorts in upper atmosphere winds → flutter and phase instability"),
        (1, "Phase instability can hurt MSK144 even when signal looks strong"),
        (1, "May decode on the leading edge before the trail saturates"),
        (0, "Both types appear in normal operation — knowing which helps interpret results"),
    ], notes="This explains why operators sometimes see a strong signal on the waterfall "
             "but no decode — overdense flutter breaks the phase coherence MSK144 needs.")

    content_slide(prs, "Ping Shape and Duration", [
        (0, "Trail formation: meteoroid ablates in ~10–100 ms"),
        (1, "Rise time is nearly instantaneous — follows the ablation event"),
        (0, "Underdense trail decay: exponential"),
        (1, "Time constant ~50–300 ms — most energy gone within half a second"),
        (1, "Trail expands by diffusion — electron density drops, signal fades"),
        (0, "Overdense trail decay: slower, more complex"),
        (1, "Can sustain signal for seconds"),
        (1, "Often shows flutter as the trail distorts in winds"),
        (0, "Total ping duration: 72 ms (minimum MSK144 decode) to several seconds"),
        (0, "MSK144 designed to decode from a single 72 ms frame"),
    ], notes="The exponential fall is why the detector captures 500 ms before "
             "and 1200 ms after the detection peak — the tail carries decodable signal.")

    content_slide(prs, "Short Underdense Pings — The Hardest Case", [
        (0, "Gamma-envelope model:  A(t) = (t/τ) · exp(−t/τ)"),
        (1, "Named for the Gamma distribution — the formula is Gamma(2, τ) in shape"),
        (1, "Rise  (t/τ):       linear — electrons deposited as meteoroid ablates"),
        (1, "Decay exp(−t/τ):   exponential — ionized trail diffuses radially outward"),
        (1, "Peak at t = τ, then exponential fall"),
        (1, "Short underdense pings: τ = 30–100 ms"),
        (0, "Implication for detection"),
        (1, "Detection hop = 42.7 ms — the metric window advances once every hop"),
        (1, "A τ=30 ms ping may produce only 1–2 hops above threshold"),
        (1, "Peak signal can fall between two hops — detected metric < instantaneous SNR"),
        (0, "Implication for decoding"),
        (1, "MSK144 requires ≥72 ms (one frame) to decode — a τ=30 ms ping barely fills it"),
        (1, "Fundamental limit: pings with τ < 40 ms may not contain enough symbols"),
        (1, "Not a system bug — the protocol itself requires a minimum burst duration"),
        (0, "τ > 40 ms at ≥−1 dB SNR: MAP144 achieves 100% decode rate"),
    ], notes="This is the physics of the problem. Short pings at threshold SNR are inherently "
             "difficult. The gamma envelope is the standard meteor trail model; "
             "τ is the time constant, not the total duration.")

    content_slide(prs, "Polarization from Faraday Rotation", [
        (0, "Transmitted signal has a fixed polarization (usually Horizontal)"),
        (0, "The ionosphere rotates the polarization plane — Faraday rotation"),
        (1, "Rotation ∝ electron density × magnetic field strength × path length"),
        (1, "Varies with time of day, season, solar activity, and path"),
        (1, "At 50 MHz: hundreds of degrees — multiple full rotations possible"),
        (1, "At 144 MHz: less rotation but still significant and unpredictable"),
        (1, "Practical result: received polarization angle is effectively random"),
        (0, "Rate of change"),
        (1, "Ionospheric conditions change slowly — seconds to minutes"),
        (1, "Polarization angle is essentially constant during a single ping (~100 ms)"),
        (1, "Angle changes from ping to ping — cannot be predicted or compensated"),
        (0, "Effect on a single linear antenna"),
        (1, "Signal arrives at an unpredictable angle relative to the antenna"),
        (1, "Worst case: 90° misalignment → near-total loss (~20 dB)"),
    ], notes="The key word is 'random' — even if you knew the rotation at one moment "
             "it changes with ionospheric conditions, time of day, and path. "
             "A horizontally polarized antenna can completely miss a signal a vertical "
             "would hear clearly, and the situation reverses on the next ping.")

    content_slide(prs, "Polarization from Meteor Trail", [
        (0, "The meteor trail adds further polarization change"),
        (1, "Overdense trails act as a conducting cylinder"),
        (1, "Reflection geometry rotates and shifts the polarization plane"),
    ], notes=" "
             " "
             " ")

    content_slide(prs, "Dual Polarization Solution", [
        (0, "Polarization arrival angle is unknown and unpredictable"),
        (0, "Solution: capture both H and V simultaneously"),
        (1, "H and V together capture all of the incoming energy"),
        (1, "No signal is lost regardless of arrival angle"),
        (0, "Software combines H and V before decoding:"),
        (1, "combined = H·cos θ + V·e^{jΔφ}·sin θ"),
        (1, "θ = polarization angle — 0° is pure H, 90° is pure V, 45° is equal mix"),
        (1, "cos θ and sin θ are amplitude weights — power sums to 1, no noise amplification"),
        (1, "e^{jΔφ} corrects the feedline and receiver phase between the two paths"),
        (0, "θ and Δφ are estimated fresh for each ping by searching for maximum SNR"),
    ], notes="The search tries a 9x8 grid of (theta, delta-phi) values and picks the "
             "combination that maximises the MSK144 detection metric. "
             "The combined signal is then passed to jt9 for decoding.")

    content_slide(prs, "More about Phase Correction — Δφ", [
        (0, "Sources of phase difference between H and V antenna paths"),
        (1, "Cable length difference — even a few cm of coax shifts phase at VHF"),
        (1, "Connector and hardware asymmetry — preamps, splitters, feed networks"),
        (1, "Antenna geometry — crossed dipoles introduce a fixed phase relationship"),
        (1, "Trail reflection — overdense trails add a ping-to-ping phase shift"),
        (0, "The hardware contribution is fixed and stable"),
        (1, "Same cables, same connectors, same offset every time"),
        (0, "The trail contribution varies ping to ping"),
        (1, "Depends on trail orientation, density, and geometry"),
        (0, "How Δφ is estimated"),
        (1, "Not measured directly — searched over 8 values: 0°, 45°, 90° … 315°"),
        (1, "Combined with the θ search: 9 angles × 8 phases = 72 combinations per ping"),
        (1, "The combination that maximises the MSK144 detection metric wins"),
        (1, "Hardware phase error and trail phase shift both corrected implicitly"),
    ], notes="The implicit search approach means no calibration procedure is needed. "
             "The system self-corrects for cable differences every ping.")

    content_slide(prs, "Doppler on the Propagation Path", [
        (0, "On 6 m (50 MHz, λ ≈ 6 m) the path applies Doppler from three sources"),
        (0, "Trail wind drift  —  upper-atmosphere wind at 90-100 km altitude"),
        (1, "Wind speed 100-200 m/s, radial Doppler ±30-65 Hz"),
        (1, "Quasi-constant within a ping; varies slowly across a session"),
        (0, "Trail expansion  —  ambipolar diffusion of the ionised column"),
        (1, "Expansion 1-5 m/s, Doppler ±0.3-1.7 Hz"),
        (1, "Monotonic decay toward zero as the trail spreads"),
        (0, "Multipath between reflection points  —  the fast component"),
        (1, "Multiple ionised regions along the trail reflect simultaneously"),
        (1, "Each at slightly different radial velocity → slightly different Doppler"),
        (1, "Their coherent sum modulates both amplitude and frequency"),
        (0, "Key insight:  amplitude flutter and frequency flutter are the same phenomenon"),
        (1, "Both are interference between paths with different Doppler"),
        (1, "A ping that flutters in amplitude is also fluttering in frequency"),
    ], notes="Set up the physics context for the coherent-integration slide that follows. "
             "Magnitudes are order-of-magnitude — the actual radial component depends on geometry. "
             "The multipath insight is what motivates the design discussion in Section 12.")

    content_slide(prs, "Coherent Integration Has a Doppler Ceiling", [
        (0, "Coherent averaging gains SNR linearly with N frames  —  IF phase is stable"),
        (1, "Signal sums as N · A (vectors add)"),
        (1, "Noise sums as √N · σ (random phases add as variance)"),
        (1, "Net gain ∝ √N when coherent — none when phase walks freely"),
        (0, "Phase must stay within ~π/4 over the full integration window"),
        (0, "Maximum allowable frequency drift per window length:"),
        (1, "1 frame  (72 ms):   ±1.70 Hz"),
        (1, "2 frames (144 ms):  ±0.87 Hz"),
        (1, "3 frames (216 ms):  ±0.58 Hz"),
        (1, "5 frames (360 ms):  ±0.35 Hz"),
        (1, "7 frames (504 ms):  ±0.25 Hz"),
        (0, "WSJT-X Deep mode (-d 3) attempts 5- and 7-frame averages"),
        (1, "Works on smooth overdense trails (one dominant reflection, stable phase)"),
        (1, "Fails silently on fluttery signals — coherent gain degrades, decoder still tries"),
        (1, "Phantom decodes during noise events come partly from this"),
    ], notes="The π/4 criterion is a standard rule of thumb for matched-filter coherent gain. "
             "Real-world measurements show this ceiling is hard to reach — see Section 12.")

    content_slide(prs, "MSK144 — The Protocol", [
        (0, "Part of the WSJT-X suite by K1JT et al."),
        (0, "Minimum Shift Keying"),
        (1, "2000 baud, 144 characters per second — designed for short bursts"),
        (1, "Decodable from fragments as short as ~72 ms"),
        (0, "Carries: callsigns, grid square, signal report"),
        (1, "13-character free-text mode also supported"),
        (0, "Calling frequencies: 50.260, 144.150, 222.100 MHz"),
        (0, "Standard tool for meteor scatter QSOs worldwide"),
    ], notes="Now placed after propagation path so the audience understands why "
             "72 ms is the design target — it matches the shortest useful underdense ping.")

    # ═══════════════════════════════════════════════════════════════════════════
    # SECTION 2 — FEATURES
    # ═══════════════════════════════════════════════════════════════════════════
    section_slide(prs, 3, "Features",
        subtitle="What MAP144 does",
        audience="Everyone",
        notes="Walk through the feature set. Use screenshots where placeholders appear.")

    content_slide(prs, "What MAP144 Does", [
        (0, "Takes a 48 kHz complex IF signal from a digital receiver"),
        (1, "Receiver tuned to the MSK144 calling frequency, looks +/- 17 kHz"),
        (1, "Supports FlexRadio, USRP B210, AirSpy, RTL-SDR, and WAV file replay"),
        (0, "Finds and decodes all MSK144 signals across ±17 kHz"),
        (1, "Creates 48 channels, on 1 kHz spacing, monitored simultaneously"),
        (1, "No operator intervention required — decodes everything it hears"),
        (0, "Shows decoded messages in a small window with color-coded aging"),
        (0, "Reports decoded messages to GridTracker, N1MM, DX cluster, PSKReporter"),
        (0, "Optional display windows:"),
        (1, "Fast Graph — real-time spectrogram and tone detection SNR display"),
        (1, "Internal decode processing — detector metrics and pipeline status"),
        (1, "Receiver status — source connection, sample rate, noise floor"),
        (1, "Oscilloscope and noise blanker — real-time IQ waveform display"),
        (0, "Analysis window — operator tool for weak or corrupted signal recovery"),
        (1, "Replay captured IQ, re-run detector, manually trigger decode attempts"),
    ], notes="One-sentence summary of the program before diving into individual features.")

    stacked_image_slide(prs, "Signal Processing Pipeline", [
        (0, "Noise blanker  —  impulse rejection before channelization"),
        (1, "Wideband FFT; hot bins zeroed and re-synthesised"),
        (0, "Polyphase channelizer  —  48-channel filter bank"),
        (1, "Each channel 1 kHz wide, decimated 48 kHz → 12 kHz complex IQ"),
        (0, "Squared-spectrum detector  —  per channel, every 42.7 ms hop"),
        (1, "MSK144 tones ±500 Hz from carrier → lobes 2000 Hz apart when squared"),
        (1, "Rolling 25th-percentile baseline over 32 s per channel"),
        (0, "Cluster-based coincidence gate  —  counts narrow-band clusters, not channels"),
        (1, "Each real ping spreads to ~3 adjacent channels; gate counts separate clusters"),
        (0, "Decimate + mix + BPF  —  12 kHz complex IQ, ±1200 Hz audio BPF"),
        (0, "SPD decoder  —  Python reimplementation of the WSJT-X Short Ping Decoder"),
        (0, "jt9 decoder  —  fallback for pings SPD misses; handles longer bursts"),
        (0, "Reporting  —  PSKReporter · WSJT-X UDP · DX cluster"),
    ], label="DSP pipeline block diagram",
       notes="Walk through each stage. Key insight: each stage addresses a specific "
             "failure mode — impulse noise → blanker; wideband → channelizer; "
             "false triggers → cluster gate; BPF eliminates detection failures; "
             "SPD decodes short pings at near-optimal sensitivity.",
       png="map144_pipeline.png")

    stacked_image_slide(prs, "Real-Time Fast Graph (like WSJT-X)", [
        (0, "48 kHz IQ stream displayed as a scrolling waterfall"),
        (1, "Top pane: current 15 s window"),
        (1, "Bottom pane: accumulated previous 15 s"),
        (0, "Frequency axis"),
        (1, "±24 kHz around the calling frequency"),
        (0, "Color scale"),
        (1, "Adjustable min/max dB"),
        (0, "Fast Graph shows test file with 10 pings"),
    ], "Main window — waterfall and spectrogram",
        notes="Show the full main window. Point out the two time windows: current 15 s and previous 15 s.",
        png="map144_fast_graph.png")

    stacked_image_slide(prs, "Tone Detection SNR Display", [
        (0, "One row per 1 kHz channel across ±24 kHz"),
        (0, "Brightness = squared-domain tone SNR above noise floor"),
        (1, "Detects MSK144 tone pair without decoding"),
        (1, "Insensitive to narrow-band interference"),
        (0, "Green circle = successful decode"),
        (0, "Orange circle = jt9 launched, no decode"),
        (0, "Time axis: 0–15 s, left to right"),
    ], "Tone Detection SNR — channel × time",
        notes="This is the key diagnostic view. Explain the channel grid and the "
              "coincidence gate that suppresses impulse noise.",
        png="map144_detection_heatmap.png")

    content_slide(prs, "Supported Radio Sources", [
        (0, "FlexRadio 6000-series"),
        (1, "FlexClient via TCP — reads IQ directly from the panadapter slice"),
        (1, "Tracks pan center frequency automatically"),
        (0, "USRP B210 (Ettus Research)"),
        (1, "Via UHD — supports dual-polarization (H and V simultaneously)"),
        (0, "AirSpy R2 / Mini"),
        (1, "Direct sampling via airspy library"),
        (0, "RTL-SDR"),
        (1, "Low-cost USB dongle — 8-bit but functional"),
        (0, "WAV file replay"),
        (1, "Offline analysis of recorded IQ — single or dual-polarization"),
    ], notes="Skip USRP if talking to a general ham audience. "
             "Emphasise FlexRadio and RTL-SDR as the most common cases.")

    bullets_image_slide(prs, "FlexRadio Status Window", [
        (0, "Discovers radios running FlexClient on the network"),
        (1, "Shows radio identity and panadapter/slice table"),
        (0, "Connects via TCP to the FlexRadio DAXIQ stream"),
        (1, "Shows DAXIQ stream status: packet rate, queue depth, loss"),
        (0, "Decodable dial frequency range — which USB freqs are covered"),
        (0, "Tracks panadapter center frequency automatically"),
        (1, "Channelizer NCO updates when pan moves"),
    ], "FlexRadio status window screenshot",
        caption="FlexRadio source window showing stream status and decodable dial frequency range.",
        notes="Key fields: Status (Receiving/Idle), pan center frequency, "
              "packet rate, queue depth, and loss counter.",
        png="map144_flex_radio.png",
        img_right_margin=0, img_align='right',
        img_top=int(Pt(8)), img_h=int(_FOOTER_TOP - Pt(8)),
        img_left=Inches(9.3), txt_w=Inches(8.9))

    bullets_image_slide(prs, "USRP B210 Status Window", [
        (0, "USRP B210 — USB 3.0 software-defined radio"),
        (1, "Originally Ettus Research, now NI (National Instruments / Emerson)"),
        (0, "Hardware identity and IF stream parameters"),
        (1, "Sample rate, center frequency, active channels"),
        (0, "Two independent RF channels — H and V for dual-polarization"),
        (1, "Independent gain slider per channel (0–76 dB)"),
        (1, "Antenna port selection: RX2 or TX-RX"),
        (1, "Signal level and noise floor readout per channel"),
        (0, "Buffer drop counter flags overrun conditions"),
    ], "USRP B210 status window screenshot",
        caption="USRP B210 source window — gain and antenna port controls for each RF channel.",
        notes="RF 0 and RF 1 are the two receive channels — H and V for dual-pol. "
              "Each has independent gain slider (0–76 dB), antenna port selection (RX2 / TX-RX), "
              "signal level, and noise floor readout. "
              "Buffer drop counter flags overrun conditions.",
        png="map144_usrp_b210.png",
        img_right_margin=0, img_align='right',
        img_top=int(Pt(8)), img_h=int(_FOOTER_TOP - Pt(8)),
        img_left=Inches(9.0), txt_w=Inches(8.6))

    content_slide(prs, "Dual-Polarization Support", [
        (0, "Antenna feeds two orthogonal polarizations: H and V"),
        (0, "Separate IQ streams captured simultaneously"),
        (0, "Independent waterfall and heatmap for each polarization"),
        (0, "Software combines H and V to maximise SNR"),
        (1, "combined = H·cos θ + V·e^{jΔφ}·sin θ"),
        (0, "Polarization angle θ estimated before each decode attempt"),
        (0, "Useful for:"),
        (1, "Faraday rotation studies"),
        (1, "Cross-polarization meteor burst experiments"),
        (1, "Antenna calibration"),
    ], notes="This is the most advanced feature. Can skip for a general audience. "
             "Good hook for the AI development section — this was developed iteratively.")

    image_slide(prs, "Estimated Polarization Angle",
        "Dual-pol analysis window showing H and V heatmaps",
        caption="θ estimated by searching a 9×8 grid of (angle, phase offset) — "
                "decoded angle shown in the decode list.",
        notes="Explain that the angle estimate feeds directly into the signal combining, "
              "improving decode probability for off-axis signals.")

    content_slide(prs, "Reporting Outputs", [
        (0, "PSKReporter  (pskreporter.info)"),
        (1, "IPFIX upload — decoded stations appear on the worldwide map"),
        (1, "Callsign, grid, frequency, SNR, mode — uploaded every 5 minutes"),
        (0, "WSJT-X UDP broadcast"),
        (1, "Emulates WSJT-X Status and Decode UDP packets on port 2237"),
        (1, "GridTracker — displays decoded stations on a map in real time"),
        (1, "N1MM Logger+ — logs contacts automatically"),
        (1, "Any software that listens for WSJT-X UDP will work"),
        (0, "DX cluster"),
        (1, "Posts spots via Telnet to your cluster node"),
        (0, "All reporting is optional and independently configurable"),
    ], notes="The WSJT-X UDP emulation is the most useful integration — "
             "GridTracker and N1MM see MAP144 decodes as if WSJT-X sent them. "
             "No changes needed in those programs.")

    image_slide(prs, "Analysis Window",
        "Analysis window — spectrogram, heatmap, audio, decode list",
        caption="Post-capture review: replay IQ, re-run detector, inspect individual pings, "
                "listen to 12 kHz audio.",
        notes="Walk through the four panes: spectrogram, detection heatmap, "
              "audio spectrogram, decode list. Show how to click a ping to decode it manually.")

    bullets_image_slide(prs, "IQ / Noise Blanker Window", [
        (0, "Time-domain plot of the received IF signal"),
        (1, "Shows impulse noise spikes clearly"),
        (1, "Aids in setting the receiver's wideband noise blanker"),
        (1, "Shows the effect of the MAP144 internal noise blanker"),
        (0, "MAP144 noise blanker — inspired by Linrad design"),
        (1, "Wideband FFT of each block; per-bin running average"),
        (1, "Block blanked if too many bins are simultaneously hot"),
        (1, "Narrowband signal (meteor burst): only a few bins elevated → passes"),
        (1, "Broadband impulse (power line click): nearly all bins hot → blanked"),
        (0, "Coincidence gate: MSK144 burst must persist across multiple hops"),
        (0, "Blanking rate shown in real time"),
    ], "IQ / Noise Blanker window screenshot",
        caption="IQ magnitude time-domain view with noise blanker controls.",
        notes="Good slide for operators near power lines or in noisy RF environments. "
              "The Linrad reference resonates with serious VHF contesters.",
        png="map144_iq_noise_blanker.png")

    # ═══════════════════════════════════════════════════════════════════════════
    # SECTION 3 — FILES GENERATED
    # ═══════════════════════════════════════════════════════════════════════════
    section_slide(prs, 4, "Files Generated",
        subtitle="What MAP144 writes to disk",
        audience="Operators and experimenters",
        notes="Useful for operators who want to post-process data or integrate with other tools.")

    content_slide(prs, "IQ Captures  —  MSK144/captures/", [
        (0, "Saved by shift-clicking the waterfall or spectrogram"),
        (0, "WAV file — raw IQ at 48 kHz"),
        (1, "Single-pol: stereo float32  (L = I,  R = Q)"),
        (1, "Dual-pol:   4-channel float32  (H_I  H_Q  V_I  V_Q)"),
        (0, "JSON sidecar — metadata"),
        (1, "Source, center frequency, sample rate, dual-pol flag"),
        (1, "Timestamp, trigger position"),
        (0, "Filename pattern:"),
        (1, "YYYYMMDD_HHMMSSZ_capture_{source}_{freq}kHz.wav"),
    ], mono_items={8},
    notes="The JSON sidecar carries everything the analysis window needs to reconstruct the context.")

    content_slide(prs, "Detection Logs  —  MSK144/detections/", [
        (0, "launches.jsonl  — one entry per jt9 launch"),
        (1, "Outcome: decoded / no_decode / timeout / error"),
        (1, "Time, frequency, estimated polarization angle θ"),
        (0, "decodes.jsonl  — one entry per successful decode"),
        (1, "Message, SNR, frequency, θ"),
        (1, "Used by the comparison report"),
        (0, "Written continuously — safe to inspect while the program runs"),
        (0, "{wavfile}_run001.txt, _run002.txt … — comparison reports"),
        (1, "Human-readable table: DECODED / NO DECODE / MISSED per signal"),
    ], mono_items={0, 3, 7},
    notes="JSONL = one JSON object per line, easy to parse with Python or jq.")

    content_slide(prs, "Decoded Audio  —  MSK144/detections/", [
        (0, "One WAV file per successful decode"),
        (0, "12 kHz mono — the audio fed to jt9"),
        (0, "Filename encodes the decode:"),
        (1, "YYYYMMDD_HHMMSSZ_{freq}kHz_{message}.wav"),
        (0, "Can be opened in Audacity, replayed through WSJT-X, or analyzed"),
        (0, "Useful for:"),
        (1, "Verifying the decode manually"),
        (1, "Comparing signal quality across antenna experiments"),
        (1, "Archive of every meteor burst heard"),
    ], mono_items={3},
    notes="These files are the 'physical evidence' of each detection. "
          "A good way to sanity-check the system.")

    # ═══════════════════════════════════════════════════════════════════════════
    # SECTION 4 — INSTALLATION
    # ═══════════════════════════════════════════════════════════════════════════
    section_slide(prs, 5, "Installation",
        subtitle="Getting MAP144 running",
        audience="Operators",
        notes="Skip this section for a general overview talk.")

    content_slide(prs, "Prerequisites", [
        (0, "Linux (primary), Windows/Mac may work with minor adjustments"),
        (0, "Python 3.11 or later"),
        (0, "jt9 binary from WSJT-X"),
        (1, "Ubuntu/Debian:  sudo apt install wsjtx"),
        (1, "Or build from source: physics.princeton.edu/pulsar/K1JT/wsjtx.html"),
        (0, "Radio-specific drivers:"),
        (1, "FlexRadio: no extra drivers — TCP connection"),
        (1, "USRP: UHD  (apt install uhd-host)"),
        (1, "AirSpy: libairspy  (apt install airspy)"),
        (1, "RTL-SDR: librtlsdr  (apt install rtl-sdr)"),
    ], mono_items={3, 4},
    notes="jt9 is the critical dependency — everything else can be worked around.")

    content_slide(prs, "Install MAP144", [
        (0, "Clone the repository:"),
        (1, "git clone https://github.com/wa1hco/map144.git"),
        (0, "Create a virtual environment:"),
        (1, "python3 -m venv .venv && source .venv/bin/activate"),
        (0, "Install Python dependencies:"),
        (1, "pip install -r requirements.txt"),
        (0, "Verify:"),
        (1, "python map144.py --help"),
    ], mono_items={1, 3, 5, 7},
    notes="Replace GitHub URL with actual repo URL before presenting.")

    content_slide(prs, "First-Run Configuration", [
        (0, "Launch MAP144:  python map144.py"),
        (0, "Select your radio source from the Source menu"),
        (0, "Set center frequency to the MSK144 calling frequency"),
        (1, "6 m: 50.260 MHz     2 m: 144.150 MHz"),
        (0, "Open Reporting window (View menu) to configure:"),
        (1, "Your callsign and grid square"),
        (1, "PSKReporter upload (on/off)"),
        (1, "WSJT-X UDP port (for GridTracker, N1MM, etc.)"),
        (1, "DX cluster hostname and port"),
        (0, "Settings are saved automatically on exit"),
    ], mono_items={0},
    notes="Walk through the menus live if possible. "
          "Settings file location: show where it lands on disk.")

    # ═══════════════════════════════════════════════════════════════════════════
    # SECTION 5 — OPERATION
    # ═══════════════════════════════════════════════════════════════════════════
    section_slide(prs, 6, "How to Operate",
        subtitle="Day-to-day use",
        audience="Operators",
        notes="This section is hands-on. A live demo works better than slides here.")

    image_slide(prs, "Starting a Live Session",
        "Source selection menu and radio parameters",
        caption="Select source → enter center frequency → press Connect. "
                "Status bar shows sample rate and source mode.",
        notes="Emphasise: center frequency must match the radio's panadapter center exactly. "
              "FlexRadio users: the app can track the pan center automatically.")

    bullets_image_slide(prs, "Reading the Waterfall", [
        (0, "Top pane: current 15 s — updates continuously"),
        (0, "Bottom pane: previous 15 s — frozen, good for review"),
        (0, "Meteor ping: brief bright vertical streak"),
        (1, "Typically 50–500 ms wide"),
        (1, "May appear on one or several adjacent channels"),
        (0, "Calling frequency shown as a horizontal reference line"),
        (0, "Color scale: adjust min/max dB to taste"),
    ], "Main window annotated: waterfall, spectrogram, noise floor",
        caption="Top: rolling 15 s.  Bottom: accumulated previous 15 s.",
        notes="Show what a meteor ping looks like — a brief bright vertical streak.",
        png="map144_main.png")

    stacked_image_slide(prs, "Reading the Tone Detection SNR Display", [
        (0, "Bright row = strong MSK144 tone pair in that channel"),
        (0, "Cyan circle  = decoded by SPD  (short ping, sensitive path)"),
        (0, "Green circle = decoded by jt9  (fallback, longer burst)"),
        (0, "Orange circle = no decode — both SPD and jt9 failed"),
        (1, "Signal too weak, too short, or too much phase distortion"),
        (0, "Red circle = decode in progress — result pending"),
        (0, "Circles appear slightly after the ping — decoder needs the full burst"),
        (0, "Detection threshold: 3 dB above 32-second per-channel baseline"),
        (1, "Baseline is 25th percentile of rolling metric history per channel"),
        (1, "32 s window is long enough to average through IF spurs"),
        (1, "Too low → false alarms from noise  |  Too high → missed weak pings"),
    ], "Tone Detection SNR annotated: channels, threshold, circles",
        notes="The threshold slider controls sensitivity. A good starting point is SNR threshold = 6 dB. "
              "Cyan vs green tells the operator which decoder succeeded — "
              "cyan (SPD) is the short-ping path optimized for underdense meteors; "
              "green (jt9) handles overdense and longer bursts.",
        png="map144_detection_heatmap.png")

    content_slide(prs, "SPD — Short Ping Decoder", [
        (0, "SPD is MAP144's primary decoder, built specifically for meteor scatter"),
        (0, "Why not just use jt9?"),
        (1, "jt9 processes a full 15-second frame — it averages over all of it"),
        (1, "A 100 ms ping is a small fraction of 15 seconds; most of the frame is noise"),
        (1, "jt9 also requires real audio; half the signal (the Q channel) is discarded"),
        (0, "SPD's key insight: coherent frame averaging"),
        (1, "One MSK144 frame = exactly 72 ms at 2000 baud"),
        (1, "Successive frames from the same ping are naturally phase-coherent"),
        (1, "Averaging N frames before decoding improves SNR by √N"),
        (1, "Operates on complex IQ — both I and Q channels, no signal loss"),
        (0, "SPD runs first on every detection; jt9 is the fallback"),
        (1, "SPD handles the common short underdense case with better sensitivity"),
        (1, "jt9 handles longer overdense bursts that span many frames"),
    ], notes="The coherent averaging gain is real: averaging 3 frames (the typical underdense ping) "
             "gives √3 ≈ 1.7× SNR improvement = +2.5 dB before any decoding. "
             "Combined with the complex IQ path (+3 dB over real audio) the SPD sensitivity "
             "advantage over raw jt9 on a short burst is approximately 5 dB.")

    content_slide(prs, "Decode Path: What the Circle Colors Mean", [
        (0, "Every detection triggers a two-stage decode attempt:"),
        (0, "Stage 1 — SPD  (Short Ping Decoder)"),
        (1, "Runs first on every detection"),
        (1, "Optimized for 1–3 frame underdense pings"),
        (1, "Operates on complex IQ with audio BPF — near-optimal sensitivity"),
        (1, "Success → cyan circle"),
        (0, "Stage 2 — jt9  (fallback, only if SPD fails)"),
        (1, "Handles longer overdense bursts and signals SPD misses"),
        (1, "Reads real audio from the 30-second IQ ring buffer"),
        (1, "Success → green circle"),
        (0, "Both fail → orange circle"),
        (0, "In practice: most short underdense pings → cyan"),
        (1, "Overdense trails, multi-second bursts → green"),
        (1, "Noise false alarms, very weak signals → orange"),
    ], notes="The two-stage path gives the best of both decoders: "
             "SPD's sensitivity for short pings, jt9's robustness for long bursts. "
             "The color tells the operator which path succeeded without reading the log.")

    content_slide(prs, "Capturing IQ for Analysis", [
        (0, "Shift-click anywhere on the waterfall or heatmap"),
        (1, "Saves 30 seconds of IQ centerd on the click time"),
        (0, "A notification confirms the save location"),
        (0, "To open for analysis:"),
        (1, "File → Open Capture…"),
        (1, "Or click the Analysis button if shown"),
        (0, "Analysis window replays the capture through the full pipeline"),
        (1, "Re-runs detection with current settings"),
        (1, "Re-runs jt9 on each detected ping"),
        (1, "Allows manual decode attempts by clicking the heatmap"),
    ], notes="The shift-click save is the primary workflow for investigating a ping. "
             "Useful when you see something on the waterfall but jt9 missed it.")

    image_slide(prs, "Analysis Window Walkthrough",
        "Analysis window — full view",
        caption="Top: IQ spectrogram.  Middle: detection heatmap.  "
                "Bottom-left: decode list.  Bottom-right: 12 kHz audio spectrogram.",
        notes="Walk through: run analysis → circles appear → click a circle → audio updates. "
              "Point out the polarization angle column in the decode list (dual-pol only).")

    content_slide(prs, "Tips and Gotchas", [
        (0, "Sample rate: 48 kHz is standard; others may need re-testing"),
        (0, "jt9 must be on PATH  (test: jt9 --version in a terminal)"),
        (0, "Noise blanker threshold: start at 6× noise floor"),
        (1, "Reduce if too much is blanked;  raise if false alarms persist"),
        (0, "PSKReporter uploads every 5 minutes — not real-time"),
        (0, "WAV replay is faster than real-time — normal behavior"),
    ], notes="These are the most common operator questions.")

    # ═══════════════════════════════════════════════════════════════════════════
    # SECTION 6 — ARCHITECTURE
    # ═══════════════════════════════════════════════════════════════════════════
    section_slide(prs, 7, "System Architecture",
        subtitle="How it works under the hood",
        audience="Technical — skip for general audience",
        notes="Go deep here for SDR/DSP audiences. Skip for general ham radio clubs.")

    image_slide(prs, "DSP Pipeline Overview",
        "Block diagram: ADC → noise blanker → channelizer → detector → SPD/jt9",
        caption="IQ stream at 48 kHz → 48-channel polyphase filter bank → "
                "per-channel squared-FFT detector → SPD decoder (jt9 fallback)",
        png="map144_pipeline.png",
        notes="Key insight: the channelizer converts a wideband IQ stream into "
              "48 narrow-band channels, one per 1 kHz slot. "
              "After the coincidence gate, complex IQ is BPF-filtered and fed to "
              "the SPD decoder first; jt9 handles the fallback for longer bursts.")

    content_slide(prs, "Polyphase channelizer", [
        (0, "48-channel filter bank centerd on the calling frequency"),
        (0, "Each channel: 1 kHz wide, decimated to 12 kHz"),
        (1, "Channel k captures signals with USB dial at calling_freq + k×1000 Hz"),
        (0, "NCO offset aligns the channel grid to the calling frequency"),
        (1, "channel_offset_hz = (calling_freq − pan_center) + 1500 Hz"),
        (0, "Output: complex IQ at 12 kHz per channel, all channels in parallel"),
        (0, "Implemented with scipy lfilter and numpy — runs in real time at 48 kHz"),
    ], mono_items={2, 4},
    notes="The +1500 Hz offset is because USB audio places signal energy 1500 Hz above "
          "the dial frequency. This aligns the NCO with the actual carrier.")

    content_slide(prs, "Detection: Squared-Spectrum Metric", [
        (0, "MSK144 tones sit ±500 Hz from the carrier — squaring maps f→2f, placing lobes 2000 Hz apart"),
        (0, "Detection metric: mean power in LO band × mean power in HI band"),
        (1, "metric = mean(|FFT|² in LO) × mean(|FFT|² in HI)"),
        (0, "Squaring makes it sensitive to simultaneous energy in both lobes"),
        (1, "Random noise: products cancel on average"),
        (1, "MSK144 burst: both bands bright simultaneously → metric spikes"),
        (0, "Threshold: 3 dB above 25th-percentile of rolling 32-second history per channel"),
        (1, "Long window ensures IF spurs don't bias the baseline"),
        (0, "Hop size: 42.7 ms — a 72 ms short ping may span only 1–2 hops"),
        (1, "Peak signal may fall between hops → detection metric lower than instantaneous SNR"),
        (1, "Gamma-envelope pings (τ = 30–100 ms) are the hardest short-underdense case"),
    ], notes="The hop timing is a fundamental limit for very short pings. "
             "A τ=30 ms ping may produce only one hop above threshold. "
             "Lowering the threshold to 3 dB (from 3.5 dB) improved short-ping detection "
             "significantly without a measurable increase in false alarm rate.")

    content_slide(prs, "Coincidence Gate — Cluster Count Design", [
        (0, "Problem: impulse noise elevates many channels simultaneously"),
        (1, "A single 1 µs impulse raises all 48 channels at once — obvious false trigger"),
        (0, "Original gate: suppress if ≥ 8 channels exceed threshold in the same hop"),
        (1, "Bug: each real MSK144 ping spreads into ~3 adjacent channels"),
        (1, "Squaring maps ±500 Hz tones to ±1000 Hz; those straddle the 1 kHz channel boundaries"),
        (1, "Three simultaneous pings at 5 kHz spacing → 3 × 3 = 9 channels → gate fires incorrectly"),
        (0, "Fix: count separate narrow-band clusters, not raw channels"),
        (1, "Cluster = contiguous run of channels; real ping = one cluster of span ≤ 3"),
        (1, "Broadband noise = one wide cluster (span > 8 channels) or many clusters simultaneously"),
        (1, "Gate fires if ≥ 6 fresh narrow clusters in the same hop"),
        (0, "Three simultaneous real pings → 3 clusters → well below the threshold of 6"),
    ], notes="The 3-channel spread is a direct consequence of the squaring operation: "
             "the ±1000 Hz tone pair straddles the 1 kHz channel spacing. "
             "The original channel-count gate was correct in principle but wrong in threshold — "
             "it did not account for normal multi-station activity.")

    content_slide(prs, "Sliding Window Intervals — Three Layers", [
        (0, "Three independent sliding windows operate at different rates and purposes"),
        (0, "Layer 1 — Channelizer tone detector  (processing.py)"),
        (1, "Window: 512 samples at 12 kHz = 42.7 ms  (one squarer block)"),
        (1, "Hop:    256 samples = 21.3 ms  (50 % overlap)"),
        (1, "Purpose: trip wire — detect tone-pair energy above 25th-pct baseline"),
        (1, "Not aligned to MSK144 frame boundaries; just looking for power"),
        (0, "Layer 2 — SPD burst search  (msk144_spd.py)"),
        (1, "Window: 864 samples = 72 ms  (one MSK144 frame at 12 kHz)"),
        (1, "Hop:    216 samples = 18 ms  (quarter-frame)"),
        (1, "Purpose: find exact burst start after channelizer fires"),
        (1, "Quarter-frame step ensures burst onset is captured within ±9 ms"),
        (0, "Layer 3 — Pre-burst SNR estimator  (detection.py)"),
        (1, "Window: 1024 samples = 85 ms  (wider for frequency resolution)"),
        (1, "Hop:    512 samples = 42.7 ms  (50 % overlap)"),
        (1, "Purpose: noise floor estimate over 500–1500 ms pre-burst segment"),
    ], notes="The quarter-frame hop in SPD (216 samples) is the smallest — it determines "
             "how precisely SPD can locate the burst onset. A half-frame hop would allow "
             "the burst to be offset by up to 36 ms from the nearest window start, "
             "which can shift the sync correlation enough to miss the sync word. "
             "The channelizer 50% hop (256 samples) is a standard overlap-add choice — "
             "finer would cost CPU without improving detection probability since the "
             "trigger threshold works on energy, not phase alignment.")

    content_slide(prs, "Decode Path: SPD First, jt9 Fallback", [
        (0, "After coincidence gate fires, decode attempt starts immediately"),
        (0, "Step 1 — Capture IQ and estimate SNR (both paths)"),
        (1, "Read [−500 ms, +1200 ms] window from ring buffer; zero-pad 100 ms each end"),
        (1, "Mix fc_hz → 1500 Hz audio centre; decimate 48→12 kHz"),
        (1, "est_snr_db: overlap-averaged PSD over 500 ms pre-burst noise — 9 frames"),
        (0, "Step 2 — Try SPD decoder on complex baseband IQ"),
        (1, "Short-Ping Decoder: optimized for 1–3 frame bursts"),
        (1, "Searches frequency, aligns sync, averages frames, runs BP+LDPC"),
        (1, "If decoded → log [MSK144 SPD] with est_snr_db, save WAV, report"),
        (0, "Step 3 — jt9 fallback if SPD returns None"),
        (1, "Re-read ring buffer: extend pre-burst to 1500 ms — 34 PSD frames"),
        (1, "Re-mix + re-decimate; recompute est_snr_db from wider noise window"),
        (1, "Write extended WAV for jt9; handles longer bursts and overdense trails"),
        (1, "If decoded → log [MSK144 DECODE] with est_snr_db, save WAV, report"),
    ], notes="SPD path stays lean — 500 ms pre-burst is sufficient because SPD works on "
             "complex IQ and doesn't need a noise reference for its own decoding. "
             "jt9 needs a noise baseline for its SNR estimate; re-reading 1500 ms gives "
             "3× more averaging (34 frames vs 9). The ring buffer always holds this history "
             "since we already waited for 1200 ms post-detection IQ to arrive. "
             "The BPF is the key insight: it narrows noise bandwidth from 4800 Hz to 2400 Hz. "
             "In the squared domain, noise×noise cross-terms drop as BW² — a factor of 4. "
             "Result: detection failures drop from 64% to 0%; system is decoder-limited.")

    image_slide(prs, "SPD Decoder Internal Pipeline",
        "SPD decoder block diagram",
        caption="Complex IQ → Audio BPF → Squared-spectrum detection → "
                "Freq search → Sync correlation → Frame averaging → BP+LDPC decode",
        png="map144_spd_detail.png",
        notes="Each stage narrows the candidate set. "
              "The audio BPF (±1200 Hz) is the critical addition from the sensitivity investigation: "
              "it eliminates all detection failures at threshold SNR. "
              "Red dashed arrows show the three exit points that route to jt9 fallback.")

    content_slide(prs, "Dual-Polarization Signal Combining", [
        (0, "H and V channels processed independently through channelizer + detector"),
        (0, "At decode time: search for θ and Δφ that maximise SNR metric"),
        (1, "combined(θ, Δφ) = H·cos θ + V·e^{jΔφ}·sin θ"),
        (0, "9×8 grid search: θ ∈ {0°,10°,…,80°},  Δφ ∈ {0°,45°,…,315°}"),
        (0, "Search performed at 12 kHz (decimated) to reduce noise bandwidth"),
        (1, "At 48 kHz noise is 4× wider — metric becomes noise-dominated"),
        (0, "Best (θ, Δφ) applied to full 48 kHz signal for jt9 decode"),
        (0, "Estimated θ stored with the decode record"),
    ], notes="The noise-bandwidth issue is a key engineering insight — "
             "and a good story for the AI development section.")

    # ═══════════════════════════════════════════════════════════════════════════
    # SECTION 7 — BUILDING WITH AI
    # ═══════════════════════════════════════════════════════════════════════════
    section_slide(prs, 8, "Building with AI",
        subtitle="A case study in AI-assisted development",
        audience="AI/software audience — optional for ham radio clubs",
        notes="This section stands alone as a talk about AI-assisted development. "
              "The project is the vehicle; the process is the story.")

    content_slide(prs, "Project Genesis", [
        (0, "Starting point: basic MSK144 concept and SDR knowledge"),
        (0, "Goal: a working real-time monitor, not a research prototype"),
        (0, "Approach: iterative development with Claude as the coding partner"),
        (0, "Not 'generate me an app' — more like pair programming"),
        (1, "Describe a feature, review the code, test, report back"),
        (1, "Catch bugs, explain what the code should do, iterate"),
        (0, "The domain knowledge stayed with the human"),
        (1, "Claude knew Python / Qt / DSP syntax"),
        (1, "The engineer knew project goals, how to test, MSK144, meteor scatter, antenna physics"),
    ], notes="Set the expectation: this is not 'AI wrote the app'. "
             "It's a collaboration where the human provided direction and domain expertise.")

    two_col_slide(prs,
        "What AI Was Good At",
        left_items=[
            (0, "Boilerplate and scaffolding"),
            (1, "Qt widget layout"),
            (1, "File I/O patterns"),
            (1, "Logging infrastructure"),
            (0, "Refactoring"),
            (1, "Rename modules without missing imports"),
            (1, "Extract a class from tangled code"),
            (0, "DSP math"),
            (1, "Filter design (scipy)"),
            (1, "FFT windowing details"),
            (1, "Polyphase filter bank structure"),
        ],
        right_items=[
            (0, "API and library usage"),
            (1, "python-pptx, pyqtgraph, UHD"),
            (1, "Correct kwargs on first try"),
            (0, "Writing tests"),
            (1, "Generating test signals at known SNR"),
            (1, "Batch performance evaluation"),
            (0, "Documentation"),
            (1, "Docstrings, module headers"),
            (1, "This slide deck"),
        ],
        left_title="Strengths", right_title="Also good",
        notes="Claude is essentially a very fast, well-read junior developer. "
              "Strong on syntax and patterns, needs guidance on intent.")

    content_slide(prs, "A Case Study: The Polarization Search", [
        (0, "Task: estimate polarization angle θ before calling jt9"),
        (0, "First implementation: search over last 512 samples of the extract window"),
        (1, "Bug 1: those 512 samples were 1.2 s after the detection — pure noise"),
        (1, "Bug 2: frequency masks built for 12 kHz applied to 48 kHz data"),
        (0, "Fix attempt 1: center the search window on the detection moment"),
        (1, "Result: went backwards — more misses than before"),
        (0, "Root cause: at 48 kHz, noise bandwidth is 4× wider"),
        (1, "Squared-spectrum metric became noise-dominated regardless of window"),
        (0, "Correct fix: decimate to 12 kHz before the pol search"),
        (1, "Same metric, correct noise bandwidth — search now finds the right angle"),
    ], notes="This is the most instructive example in the project. "
             "The first 'fix' made things worse because the root cause wasn't understood. "
             "Claude generated the code; the engineer had to diagnose why it failed.")

    content_slide(prs, "What Required Domain Expertise", [
        (0, "Knowing that the metric was designed for 12 kHz data"),
        (0, "Understanding why 48 kHz breaks the noise assumption"),
        (0, "The +1500 Hz USB audio convention — not obvious from code alone"),
        (0, "MSK144 spectral shape: two lobes, why the product metric works"),
        (0, "Recognising when a 'fix' was moving in the wrong direction"),
        (1, "The regression was subtle — more misses, not a crash"),
        (0, "Deciding what to test and what 'good' looks like"),
        (1, "Generated test signals at known SNR and angle"),
        (1, "Batch analysis across many signals to measure detection rate"),
    ], notes="The key message: AI accelerated implementation; "
             "human expertise was essential for correctness and diagnosis.")

    content_slide(prs, "Development Statistics", [
        (0, "~62 commits over the development period"),
        (0, "~15 000 lines of Python across 22 application modules"),
        (0, "Features added iteratively: detection → display → reporting → dual-pol → sensitivity"),
        (0, "Time: significantly faster than solo development"),
        (1, "Estimate: 3–5× speed-up on implementation tasks"),
        (1, "Debugging time similar — AI cannot test against hardware"),
        (0, "Cost: API usage — pennies per session"),
        (0, "No prior experience with PyQt5, pyqtgraph, or polyphase filter banks"),
        (0, "Cross-session continuity: CLAUDE.md + persistent memory files"),
        (1, "Findings, thresholds, and architectural decisions survive session boundaries"),
    ], notes="Run: git log --oneline | wc -l    and    wc -l map144_app/*.py "
             "to get current numbers before presenting.")

    content_slide(prs, "Lessons Learned", [
        (0, "AI is a force multiplier, not a replacement for expertise"),
        (0, "The conversation is the product — describe intent clearly"),
        (1, "Vague prompt → plausible-looking but wrong code"),
        (1, "Precise description of expected behavior → good code"),
        (0, "Test everything, especially at system boundaries"),
        (1, "AI-generated DSP code can be mathematically correct but contextually wrong"),
        (0, "Regressions happen — maintain a reproducible test signal"),
        (0, "AI has no ego — it will rewrite anything without complaint"),
        (1, "This makes iterative refinement very fast"),
        (0, "Good for ham radio projects: moderate complexity, narrow domain"),
    ], notes="These lessons are transferable to any AI-assisted project. "
             "The ham radio domain is a good fit: well-defined protocols, "
             "measurable outcomes, and the developer has deep domain knowledge.")

    # ── Performance: Python + numba ───────────────────────────────────────────

    content_slide(prs, "A Case Study: Finding the CPU Bottleneck", [
        (0, "Tool: py-spy — statistical profiler, attaches to a running process"),
        (1, "sudo py-spy top --pid $(pgrep -u $USER -f map144.py)"),
        (0, "Test: 50-ping synthetic file, SNR −1 to +2 dB (all near-threshold)"),
        (0, "Result: one function consumed 97% of all CPU"),
        (1, "_bpdecode128_90  —  137 seconds for 50 pings"),
        (1, "A 15-second file took 165 seconds to fully process"),
        (1, "Everything else — channelizer, display, radio I/O — was rounding error"),
        (0, "Root cause: pure-Python nested loops over the LDPC parity check graph"),
        (1, "BP iterations: 10 max per frame, each iterates over N=128 variable nodes"),
        (1, "Python executes each loop step as interpreted bytecode — inherently slow"),
        (1, "No numpy vectorisation possible — the graph structure defeats array operations"),
    ], mono_items={1},
    notes="The profiler made the diagnosis unambiguous. Without it, this bottleneck "
          "could have been invisible — the program works correctly, just slowly. "
          "py-spy attaches to a live process with no code changes required.")

    # BP profiling table
    _bp_slide = _title_content(prs)
    add_rect(_bp_slide, 0, 0, W, Pt(8), ACCENT)
    _fit_placeholders(_bp_slide)
    _set_title_placeholder(_bp_slide, "py-spy Profile: 50-Ping Test File")
    _fill_body_placeholder(_bp_slide, [
        (0, "Top functions by cumulative CPU time (137 s wall clock, 50 pings):"),
    ])
    _add_data_table(_bp_slide,
        headers=["Function", "File", "%Own", "Own Time"],
        rows=[
            ["_bpdecode128_90",    "msk144_decode.py",  "97%", "137.3 s"],
            ["_sync_correlate",    "msk144_spd.py",     " 0%",   "4.0 s"],
            ["_freq_search_avg",   "msk144_spd.py",     " 1%",   "2.1 s"],
            ["_rescaleData_nditer","pyqtgraph",          " 1%",   "1.7 s"],
            ["_fir_filt_2d",       "channelizer.py",    " 0%",   "0.5 s"],
            ["Everything else",    "—",                 " 1%",   "1.5 s"],
        ],
        col_widths_in=[3.8, 3.5, 1.2, 1.5],
        left=Inches(0.9), top=Inches(2.5), row_h_in=0.46)
    add_textbox(_bp_slide,
        "One function.  97%.  No ambiguity about where to focus.",
        Inches(0.4), Inches(5.65), Inches(12.5), Inches(0.6),
        color=ACCENT, size=20, bold=True)
    add_notes(_bp_slide,
        "The 107% Active / 89% GIL numbers confirm this is pure Python compute — "
        "the GIL is held almost continuously because there is no I/O or C extension "
        "releasing it. The profiler makes the fix obvious.")

    content_slide(prs, "The Fix: One Decorator with Numba JIT", [
        (0, "numba.njit compiles Python functions to machine code via LLVM"),
        (1, "Same source code — add @numba.njit(cache=True) and it compiles"),
        (1, "First call: ~1 second to JIT-compile; cached to disk afterward"),
        (1, "All subsequent calls: execute at near-C speed"),
        (0, "Requirements for numba compatibility:"),
        (1, "Numeric types only — int, float, numpy arrays"),
        (1, "No Python objects, no Optional returns, no bytearray"),
        (0, "Solution: extract the inner loops into _bp_inner (numba)"),
        (1, "CRC check stays in Python — called at most once per frame"),
        (1, "Outer _bpdecode128_90 wraps _bp_inner, handles Optional return"),
        (0, "Pure-Python fallback preserved — numba is optional at import time"),
        (1, "if _NUMBA: use _bp_inner   else: original Python loops"),
    ], mono_items={1, 3, 8, 10, 11},
    notes="The cache=True is important — without it, every map144 launch pays "
          "the JIT compile cost. With it, the compiled code is stored in __pycache__ "
          "and reused across restarts.")

    # Benchmark table
    _bench_slide = _title_content(prs)
    add_rect(_bench_slide, 0, 0, W, Pt(8), ACCENT)
    _fit_placeholders(_bench_slide)
    _set_title_placeholder(_bench_slide, "Numba Result: 2300× Speedup")
    _fill_body_placeholder(_bench_slide, [
        (0, "Benchmark: 1000 calls to _bpdecode128_90, threshold-SNR LLRs:"),
    ])
    _add_data_table(_bench_slide,
        headers=["Implementation", "1000 calls", "Per call", "50-ping file"],
        rows=[
            ["Pure Python",      "~137 s (est.)", "~137 ms",  "165 s backlog"],
            ["Numba @njit",      "0.060 s",        "0.06 ms",  "< 1 s"],
            ["Speedup",          "~2300×",          "~2300×",  "real-time"],
        ],
        col_widths_in=[3.5, 2.8, 2.2, 3.0],
        left=Inches(0.9), top=Inches(2.5), row_h_in=0.50)
    add_textbox(_bench_slide,
        "The 15-second file that created a 165-second processing backlog now clears in real time. "
        "No algorithm change — identical numerical results. Just compiled loops.",
        Inches(0.4), Inches(4.8), Inches(12.5), Inches(0.9),
        color=MUTED, size=18, italic=True)
    add_notes(_bench_slide,
        "The pure Python estimate is extrapolated from the profiling data. "
        "The numba number is a direct measurement. "
        "2300× speedup from a single decorator — this is the best-case scenario "
        "for numba: a tight numeric loop with no Python objects and no I/O.")

    content_slide(prs, "Python for DSP: The Complete Picture", [
        (0, "Common concern: 'Python is too slow for real-time DSP'"),
        (0, "The reality is more nuanced — Python has three speed tiers:"),
        (1, "Tier 1 — Pure Python loops:  slow  (bytecode, GIL, overhead per op)"),
        (1, "Tier 2 — NumPy array ops:    fast   (C/BLAS under the hood, batch ops)"),
        (1, "Tier 3 — Numba @njit:        very fast  (LLVM machine code, ~C speed)"),
        (0, "MAP144 uses all three appropriately:"),
        (1, "Channelizer, BPF, FFT, detection metric → NumPy (already fast)"),
        (1, "BP decoder inner loops → Numba (Python syntax, C performance)"),
        (1, "UI, logging, config, reporting → pure Python (speed irrelevant)"),
        (0, "What about Fortran?"),
        (1, "Fortran BP decoder verified 100% identical output to Python"),
        (1, "Numba BP decoder is now equally fast — Fortran not needed"),
        (1, "Python + NumPy + Numba covers the full performance range"),
    ], notes="This is the key takeaway for software-defined radio work in Python. "
             "You don't need C or Fortran for DSP performance if you use the right tool "
             "for each layer. The Fortran investigation was still valuable — it proved "
             "the algorithm was correct before optimising it.")

    # ═══════════════════════════════════════════════════════════════════════════
    # SECTION 9 — PERFORMANCE ENGINEERING
    # ═══════════════════════════════════════════════════════════════════════════
    section_slide(prs, 9, "Performance Engineering",
        subtitle="Iterative profiling and optimization — five rounds",
        audience="Technical",
        notes="This section traces the full profiling campaign from first profile to "
              "the current optimised state. Each round shows what the profiler revealed, "
              "the technique used to fix it, and how the fix was verified. "
              "Good standalone talk for a software or SDR audience.")

    content_slide(prs, "What Real-Time Means Here", [
        (0, "MAP144 must keep up with dual 48 kHz IQ stream — continuously"),
        (0, "Core work done every 20 ms (one receive block, 960 output samples):"),
        (1, "Channelizer: 48-channel polyphase FIR bank → (48, 960) complex64"),
        (1, "Squarer + detection metric: per-channel SNR, threshold, cluster gate"),
        (1, "Ring buffer write, IQ capture for jt9 launch"),
        (1, "Display update: waterfall FFT, spectrogram, detection plots"),
        (0, "When a ping is detected (every few seconds during a shower):"),
        (1, "IQ extract → decimate → SPD decode attempt"),
        (1, "jt9 subprocess launched with 12 kHz audio"),
        (0, "Budget: all of the above on a single laptop CPU, plus headroom for jt9"),
        (0, "When any stage exceeds its slot, a backlog grows — missed pings, stale display"),
    ], notes="Set the real-time constraint explicitly before the profiling story. "
             "The audience needs to understand that 'slow' means lost decodes, not just "
             "a sluggish UI.")

    content_slide(prs, "The Profiling Tool: py-spy", [
        (0, "py-spy — statistical sampling profiler for Python"),
        (1, "Attaches to a running process with no code changes"),
        (1, "sudo py-spy top --pid $(pgrep -u $USER -f map144.py)"),
        (1, "Samples the call stack at ~100 Hz — no instrumentation overhead"),
        (0, "Key columns in the output:"),
        (1, "Own%    — time spent inside the function body (excludes callees)"),
        (1, "Total%  — time in this function including all callees"),
        (1, "OwnTime — absolute seconds in the function body over the profiling window"),
        (0, "Two headline metrics reveal the bottleneck type:"),
        (1, "Active%: fraction of time at least one Python frame was on-CPU"),
        (1, "GIL%:    fraction of time the GIL was held — high means Python bytecode"),
        (1, "GIL% ≈ Active% → pure Python compute; GIL% ≪ Active% → C extension work"),
        (0, "Workflow: run under realistic load, sample for 60–120 s, compare across rounds"),
    ], mono_items={1},
    notes="The sudo is needed to attach to a process owned by the same user on some "
          "Linux configurations. The GIL/Active relationship is the key diagnostic: "
          "if both are high, the fix is JIT or vectorisation. "
          "If Active is high but GIL is low, the bottleneck is already in C but doing "
          "too much work — restructure or use a better API.")

    # ── Round 1 table ─────────────────────────────────────────────────────────
    _r1_slide = _title_content(prs)
    add_rect(_r1_slide, 0, 0, W, Pt(8), ACCENT)
    _fit_placeholders(_r1_slide)
    _set_title_placeholder(_r1_slide, "Round 1 — The BP Decoder (97% of CPU)")
    _fill_body_placeholder(_r1_slide, [
        (0, "First profiling run.  One function.  Unmistakable."),
        (0, "py-spy profile: 137 s wall clock, 50-ping synthetic test file"),
    ])
    _add_data_table(_r1_slide,
        headers=["Function", "Own%", "OwnTime", "What it does"],
        rows=[
            ["_bpdecode128_90",  "97%", "133.1 s", "LDPC belief-propagation decoder"],
            ["_sync_correlate",  " 0%",   "4.0 s", "MSK144 sync correlation"],
            ["_freq_search_avg", " 1%",   "2.1 s", "Frequency search"],
            ["Everything else",  " 2%",   "2.8 s", "Channelizer, display, I/O"],
        ],
        col_widths_in=[3.6, 1.0, 1.8, 5.3],
        left=Inches(0.9), top=Inches(2.55), row_h_in=0.46)
    add_textbox(_r1_slide,
        "Fix: @numba.njit(cache=True) on the inner BP loop.  "
        "Result: 137 ms → 0.06 ms per decode — 2300×.  "
        "GIL% fell from 89% to single digits.  A 165 s backlog cleared to real-time.",
        Inches(0.4), Inches(5.35), Inches(12.5), Inches(1.0),
        color=MUTED, size=17, italic=True)
    add_notes(_r1_slide,
        "The full story is in Section 8. Summarised here as round 1 of the profiling "
        "campaign so the table at the end of this section has a complete baseline.")

    content_slide(prs, "Round 2 — NumPy Dispatch Overhead (_wrapfunc)", [
        (0, "After Numba fix, a new entry appeared: _wrapfunc at 42 s (5% own)"),
        (0, "Root cause: NumPy function-form calls route through a Python dispatcher"),
        (1, "np.sort(arr)      → _wrapfunc → arr.__array_function__ → C sort"),
        (1, "arr.sort()         → direct C method call — no Python dispatch"),
        (0, "The dispatch layer adds one full Python stack frame per call"),
        (1, "Called at 50 Hz × many channels — overhead accumulates"),
        (0, "Affected calls in the hot path:"),
        (1, "np.sort(...)         → arr.sort()                     (clustering)"),
        (1, "np.partition(arr, k) → arr.partition(k)               (SNR percentile)"),
        (1, "np.median(arr)       → arr.partition(k); return arr[k] (noise floor)"),
        (0, "The np.median fix is the most instructive:"),
        (1, "np.median → np.quantile → _quantile → partition — four levels deep"),
        (1, "Replaced with: copy → .partition(k) → index  — direct, in-place, O(n)"),
        (0, "Rule: prefer method forms (.sort, .partition) over function forms in hot loops"),
    ], mono_items={1, 2, 7, 8, 9, 11, 12},
    notes="This is a subtle but measurable effect. The _wrapfunc mechanism exists to "
          "support array protocol dispatch (dask, cupy, etc.). Fine at the outer edge "
          "of the program; significant overhead when called thousands of times per second.")

    content_slide(prs, "Round 3 — Channelizer: Heap Allocations and the NCO", [
        (0, "Profile after Round 2: apply_channelizer 30.6 s, _fir_filt_2d_c64 14.7 s"),
        (0, "Old channelizer path per 20 ms block:"),
        (1, "Phase matrix (n_ch × N): n_ch × N complex64 — new heap alloc"),
        (1, "Mixed signal (n_ch × N): element-wise multiply raw × phase — second alloc"),
        (1, "FIR filter over (n_ch × N) → decimate → HP filter"),
        (0, "Two large allocations flushed the L2 cache on every hop — bottleneck was memory"),
        (0, "Fix: fused NCO+FIR Numba prange kernel"),
        (1, "One Numba prange loop over n_ch channels (runs parallel across cores)"),
        (1, "Per thread: local buffer of size ntaps−1+N — NCO mix written inline"),
        (1, "FIR dot product accumulated from local buffer — never writes a (n_ch, N) array"),
        (1, "Phase returned as a scalar per channel — no np.angle / np.exp upcast"),
        (0, "Single-pol:  0.591 ms → 0.405 ms  (−31%)"),
        (0, "Dual-pol:    0.835 ms → 0.672 ms  (−19%)"),
    ], mono_items={1, 2, 7, 8, 9, 10},
    notes="The insight is that the bottleneck was not arithmetic — it was the memory "
          "allocator and cache pressure from two large intermediate arrays. "
          "The fused kernel eliminates both by staging in a thread-local buffer "
          "that fits in L1 cache.")

    code_slide(prs,
        "Fused NCO+FIR Kernel — Core Loop",
        """\
@numba.njit(parallel=True, cache=True)
def _nb_mix_fir_sp(b_rev, raw, mix_phase, mix_step, zi, y, new_zi, new_mix_phase):
    ntaps = b_rev.shape[0]; N = raw.shape[0]; n_ch = mix_phase.shape[0]
    for k in numba.prange(n_ch):          # parallel over channels
        ph = mix_phase[k]; st = mix_step[k]
        buf = np.empty(ntaps - 1 + N, dtype=np.complex64)  # thread-local
        for j in range(ntaps - 1):        # copy filter state into local buf
            buf[j] = zi[k, j]
        for n in range(N):                # NCO mix: multiply inline → buf
            buf[ntaps - 1 + n] = raw[n] * ph
            ph *= st
        for n in range(N):                # FIR: dot product over pre-rev taps
            acc = np.complex64(0.0)
            for j in range(ntaps):
                acc += b_rev[j] * buf[n + j]
            y[k, n] = acc
        for j in range(ntaps - 1):        # save filter state for next block
            new_zi[k, j] = buf[N + j]
        new_mix_phase[k] = ph             # return phase — no np.exp upcast""",
        caption=(
            "b_rev: pre-reversed taps stored in ChannelizerState (no copy per call).  "
            "Thread-local buf fits in L1 cache.  "
            "NCO advance: ph *= st  (complex multiply) — no angle/exp."
        ),
        notes="The key insight: by fusing the NCO mix into the same pass as the FIR "
              "filter state load, we never materialise the (n_ch, N) phase matrix or "
              "the (n_ch, N) mixed-signal array. Both were flushing the L2 cache on "
              "every 20 ms hop. The thread-local buf is small enough to stay in L1.")

    content_slide(prs, "Round 4 — Display Windows Dominate CPU When Visible", [
        (0, "Observation: closing the Fast Graph, Detection, and IQ windows changed the profile"),
        (1, "Active% fell from 64% to 19% — two-thirds of CPU was pure display work"),
        (0, "Root cause: display updates ran unconditionally, every 20 ms block"),
        (1, "Waterfall FFT: np.fft.fft over fft_size samples, fftshift, log10 — every hop"),
        (1, "Spectrogram setImage: pyqtgraph buffer copy + render — every hop"),
        (1, "pyqtgraph's _rescaleData appeared in every profile — pure display overhead"),
        (0, "Fix: visibility gate — check window.isVisible() once per block"),
        (1, "Skip waterfall FFT and spectrogram writes when window is hidden"),
        (1, "Buffer slide (timing state) always runs — just the compute is gated"),
        (0, "Effect when all display windows open: Active ~60%, GIL ~15%"),
        (0, "Effect when all display windows closed: Active ~19%, GIL ~7%"),
        (0, "Lesson: display rendering is not free — gate it when nothing is watching"),
    ], mono_items={1, 2, 3},
    notes="This is an architectural point, not a numerical optimization. "
          "The fix is a single isVisible() check — one line of code — but the result "
          "is a 3× reduction in CPU when windows are closed. "
          "Useful for headless or overnight operation where no one is watching the display.")

    content_slide(prs, "Round 5 — The Hidden lfilter Slow Path (USRP Decimator)", [
        (0, "Profile showed: _apply (usrp_source.py) 8.98 s, apply_along_axis 6.85 s"),
        (0, "The USRP decimator: 4× FIR anti-alias before 48 kHz output"),
        (1, "192 kHz → 48 kHz, 65-tap FIR, 3840 samples per block"),
        (0, "Code used:  scipy.signal.lfilter(b, 1.0, x)"),
        (1, "a = 1.0 signals a pure FIR (no IIR denominator)"),
        (1, "In recent scipy, lfilter with a=[1] falls back to:"),
        (1, "   apply_along_axis(lambda y: np.convolve(b, y), axis, x)"),
        (1, "A Python-level loop — the 'mystery convolve' visible in every earlier profile"),
        (0, "Fix: scipy.signal.upfirdn  — dedicated FIR + decimation Cython path"),
        (1, "Processes real and imaginary parts separately as float32"),
        (1, "Streaming state: prepend 64-sample history, slice to skip warmup + tail"),
        (1, "Verified against lfilter ground truth: 9×10⁻⁷ relative error (float32 limit)"),
        (0, "Result: 0.215 ms → 0.123 ms per block  (−43%)"),
    ], mono_items={1, 3, 4, 5, 6, 8, 9, 10},
    notes="The lfilter fallback is a scipy implementation detail that changed in a "
          "relatively recent version. The behaviour is correct — it produces the right "
          "output — but it dispatches through Python for every column of the array. "
          "upfirdn is the canonical scipy function for FIR+decimation and avoids this.")

    # ── Cumulative improvement scorecard ─────────────────────────────────────
    _score_slide = _title_content(prs)
    add_rect(_score_slide, 0, 0, W, Pt(8), ACCENT)
    _fit_placeholders(_score_slide)
    _set_title_placeholder(_score_slide, "Five Rounds — Cumulative Scorecard")
    _fill_body_placeholder(_score_slide, [
        (0, "Each round: profile → identify → change → verify → repeat"),
    ])
    _add_data_table(_score_slide,
        headers=["Round", "Component", "Before", "After", "Improvement", "Technique"],
        rows=[
            ["1", "BP decoder",           "137 ms/decode",   "0.06 ms",  "2300×",   "Numba @njit"],
            ["2", "NumPy dispatch",        "42 s profiler",   "~0 s",     "eliminated", "Method forms"],
            ["3", "Channelizer (1-pol)",   "0.591 ms/block",  "0.405 ms", "−31%",    "Numba prange\nfused kernel"],
            ["3", "Channelizer (2-pol)",   "0.835 ms/block",  "0.672 ms", "−19%",    "Numba prange\nfused kernel"],
            ["4", "Display rendering",     "Active 64%",      "Active 19%","−70%",   "Visibility\ngating"],
            ["5", "USRP decimator",        "0.215 ms/block",  "0.123 ms", "−43%",    "scipy upfirdn"],
        ],
        col_widths_in=[0.8, 2.5, 2.2, 1.8, 1.7, 2.3],
        left=Inches(0.5), top=Inches(2.4), row_h_in=0.52)
    add_textbox(_score_slide,
        "Starting point: 165 s backlog on a 15 s file.  "
        "End state: real-time at ≈19% Active CPU (windows closed), "
        "≈60% Active (all displays open) — full headroom for jt9 and dual-pol.",
        Inches(0.4), Inches(6.05), Inches(12.5), Inches(0.9),
        color=MUTED, size=17, italic=True)
    add_notes(_score_slide,
        "The improvements compound: fixing round 1 exposed round 2 and 3; "
        "fixing round 3 exposed round 4 and 5. Without profiling, most of these "
        "would have been invisible — especially the lfilter slow path and the "
        "display rendering cost.")

    content_slide(prs, "Techniques Summary — Right Tool for Each Problem", [
        (0, "Numba @njit  —  pure Python numeric loops that can't be vectorised"),
        (1, "BP decoder inner loops: graph structure defeats array operations"),
        (1, "Add one decorator; identical source; ~C speed from LLVM machine code"),
        (0, "Numba @njit prange  —  parallel work with per-thread local state"),
        (1, "Channelizer: 48 independent channels, each needs its own NCO phase"),
        (1, "Thread-local buffer avoids cache conflicts; eliminates (n_ch, N) heap allocs"),
        (0, "NumPy method forms  —  avoid the array-protocol dispatch layer"),
        (1, ".sort(), .partition(k)  instead of np.sort(), np.partition()"),
        (1, "Matters at > 100 calls/second; invisible in occasional use"),
        (0, "scipy upfirdn  —  dedicated Cython path for FIR + decimation"),
        (1, "Faster than lfilter for FIR (no IIR denominator, no Python dispatch)"),
        (1, "Streaming state via prepend-and-slice — no stateful API required"),
        (0, "Architecture  —  visibility gating"),
        (1, "Don't compute what nobody sees — check isVisible() once per block"),
        (1, "Two-thirds of CPU saved when display windows are hidden"),
    ], mono_items={1, 3, 5, 7, 9, 11, 13},
    notes="The common thread: always profile first, then match the technique to the "
          "root cause. Using Numba for everything would not have fixed the display "
          "rendering or the dispatch overhead.")

    content_slide(prs, "Lessons from Five Rounds of Profiling", [
        (0, "Profile before optimising — the bottleneck is almost never where you expect"),
        (1, "Round 1 looked like a display problem; it was the decoder"),
        (1, "Round 5 looked like an FFT problem; it was a scipy API regression"),
        (0, "Read GIL% and Active% together"),
        (1, "Both high → Python bytecode compute → Numba or vectorise"),
        (1, "Active high, GIL low → already in C → restructure, better API, less allocation"),
        (0, "One fix unmasks the next — profiling is iterative by nature"),
        (0, "Verify every change numerically, not just by feel"),
        (1, "Compare output to a known-good reference: lfilter, Fortran BP, synthetic signals"),
        (1, "A faster but wrong result is worse than a slow correct one"),
        (0, "Gates and guards are free — computing invisible things is not"),
        (0, "The profiler is the most important tool in the SDR performance toolkit"),
        (1, "py-spy attaches with no code changes — use it early and often"),
    ], notes="Wrap the section with actionable takeaways. "
             "The audience should leave knowing: profile first, fix the measured "
             "bottleneck, verify, repeat. The specific techniques (Numba, upfirdn, "
             "method forms) are secondary to the discipline.")

    # ═══════════════════════════════════════════════════════════════════════════
    # SECTION 10 — DECODER SENSITIVITY INVESTIGATION
    # ═══════════════════════════════════════════════════════════════════════════
    section_slide(prs, 10, "Decoder Sensitivity Investigation",
        subtitle="A systematic hunt for 5 missing decibels",
        audience="Technical",
        notes="This section describes a days-long, multi-tool investigation carried out "
              "with Claude as the primary coding and analysis partner. "
              "The story: start with a measurable gap, decompose it into candidates, "
              "test each one rigorously, close each gap in turn.")

    content_slide(prs, "The Problem: A 5 dB Sensitivity Gap", [
        (0, "Starting point: MAP144 SPD was decoding at +4.0 dB SNR (50% threshold)"),
        (0, "Reference: WSJT-X / QEX paper reports −1.4 dB for the same mode"),
        (0, "Gap: 5.4 dB — roughly 3.5× more signal power required"),
        (0, "At the margin this matters:"),
        (1, "A ping that WSJT-X would decode, MAP144 misses"),
        (1, "Short underdense pings are right at the threshold"),
        (0, "Three candidate explanations identified:"),
        (1, "Signal path: does converting to real audio double the noise?"),
        (1, "Decoder: is the Python BP decoder less sensitive than the Fortran?"),
        (1, "SPD pipeline: overhead from detection, sync, and frame averaging?"),
    ], notes="Frame this as an engineering investigation, not a bug report. "
             "The goal is to quantify each contribution separately so fixes are targeted.")

    content_slide(prs, "The Test Framework", [
        (0, "Required: reproducible sensitivity measurements at known SNR"),
        (0, "decode_sensitivity.py — single-frame SNR sweep"),
        (1, "Generates synthetic MSK144 at 12 kHz complex baseband"),
        (1, "Adds AWGN at calibrated SNR levels (2500 Hz BW reference)"),
        (1, "Calls msk144decodeframe directly — no pipeline, no jt9 subprocess"),
        (1, "Measures 50% decode threshold across N trials per SNR level"),
        (0, "gap3_diagnose.py — full SPD pipeline diagnostic"),
        (1, "Runs the complete SPD path (detect → sync → frame avg → decode)"),
        (1, "Oracle modes: bypass detection, bypass sync to isolate each stage"),
        (1, "Reports det_miss / sync_miss / oracle_decode_miss per SNR level"),
        (0, "Fortran test executables — cross-language decoder verification"),
        (1, "Built bp_test and mf_test from the WSJT-X Fortran source"),
        (1, "Fed identical LLRs to both Python and Fortran, compared frame-by-frame"),
    ], notes="The oracle concept is the key methodological innovation: "
             "bypass one stage at a time to measure its contribution to failures. "
             "If oracle_decode% ≈ SPD%, the bottleneck is the decoder, not detection.")

    # Gap 1 ── Hilbert / real-audio path
    content_slide(prs, "Gap 1: The Real-Audio Path Doubles Noise", [
        (0, "Current pipeline: complex IQ → take real part → Hilbert transform"),
        (1, "Goal: recover the analytic signal for the sync correlator"),
        (0, "The problem: taking Re{z} discards the Q channel — half the signal"),
        (1, "Then Hilbert re-synthesizes Q from I, but includes the noise in I twice"),
        (1, "Result: noise power doubles relative to signal power — a 3 dB penalty"),
        (0, "The fix: pass complex IQ directly, skip Hilbert entirely"),
        (1, "Channelizer already produces complex baseband IQ"),
        (1, "Add complex_baseband parameter to msk144_spd_decode()"),
        (1, "Scale by √2 to match the sync correlator's amplitude convention"),
        (0, "The √2 scale matters — correlator threshold 1.3 was calibrated for the Hilbert path"),
        (1, "Without it: complex path is √2 under-scaled → 3 dB detection loss"),
        (1, "This partially cancels the 3 dB noise gain — net improvement only 0.8 dB"),
        (1, "With √2 fix: full benefit recovered"),
    ], notes="The √2 subtlety is the key lesson here — a 'correct' fix that was still "
             "wrong because of an undocumented amplitude convention in the Fortran code. "
             "Finding it required tracing the sync correlator threshold back to its calibration.")

    # Gap 1 results table
    _g1_slide = _title_content(prs)
    add_rect(_g1_slide, 0, 0, W, Pt(8), ACCENT)
    _fit_placeholders(_g1_slide)
    _set_title_placeholder(_g1_slide, "Gap 1: Results — Complex IQ Path")
    _fill_body_placeholder(_g1_slide, [
        (0, "Measured SPD 50% decode threshold vs path:"),
    ])
    _add_data_table(_g1_slide,
        headers=["Signal path", "SPD 50% threshold", "Gain vs baseline"],
        rows=[
            ["Real audio (baseline)",             "+4.0 dB", "—"],
            ["Complex IQ, no √2 fix",             "+3.2 dB", "+0.8 dB"],
            ["Complex IQ + √2 scale (Gap 1 fix)", "+1.4 dB", "+2.6 dB"],
        ],
        col_widths_in=[5.5, 3.5, 3.0],
        left=Inches(0.9), top=Inches(2.6), row_h_in=0.50)
    add_textbox(_g1_slide,
        "The √2 fix recovered 1.8 dB that a naive implementation left on the table. "
        "Remaining gap after Gap 1: +1.4 dB − (−1.4 dB reference) = 2.8 dB.",
        Inches(0.4), Inches(4.8), Inches(12.5), Inches(0.9),
        color=MUTED, size=18, italic=True)
    add_notes(_g1_slide,
        "The first implementation gained only 0.8 dB — looked like a partial win. "
        "Only by tracing the Fortran sync threshold calibration did we find the √2 error. "
        "Without the memory system this finding would have been lost between sessions.")

    # Gap 2 ── Python vs Fortran decoder
    content_slide(prs, "Gap 2: Is the Python Decoder Less Sensitive?", [
        (0, "Hypothesis: Python port of BP decoder might be less sensitive than Fortran"),
        (1, "+1.4 dB remaining gap — could the Python decoder account for it?"),
        (0, "Test strategy: build Fortran test executables from WSJT-X source"),
        (1, "Extracted bpdecode128_90.f90 and msk144decodeframe.f90"),
        (1, "Compiled with gfortran, linked against the same tables as WSJT-X"),
        (1, "Wrote a C bridge to call Fortran from Python test harness"),
        (0, "Comparison: feed identical LLRs to Python and Fortran, compare output frame-by-frame"),
        (1, "bp_test: Python bpdecode128_90 vs Fortran bpdecode128_90"),
        (1, "mf_test: full Python msk144decodeframe vs Fortran equivalent"),
        (0, "200 trials per SNR level, −2 to +2 dB range"),
    ], notes="Building the Fortran test executables was the most complex single piece of "
             "this investigation. It required understanding the WSJT-X build system, "
             "extracting the right source files, and writing interface code. "
             "Claude handled all of this from a description of what was needed.")

    # Gap 2 results table
    _g2_slide = _title_content(prs)
    add_rect(_g2_slide, 0, 0, W, Pt(8), ACCENT)
    _fit_placeholders(_g2_slide)
    _set_title_placeholder(_g2_slide, "Gap 2: Python vs Fortran — The Verdict")
    _fill_body_placeholder(_g2_slide, [
        (0, "Frame-by-frame agreement across 200 trials × 5 SNR levels:"),
    ])
    _add_data_table(_g2_slide,
        headers=["Test", "Decoder A", "Decoder B", "Agreement", "50% threshold"],
        rows=[
            ["BP decoder only",    "Python bpdecode128_90", "Fortran bpdecode128_90", "100% every frame", "—"],
            ["Full frame decoder", "Python msk144decodeframe", "Fortran mf_test", "100% every frame", "0.0 dB"],
        ],
        col_widths_in=[2.2, 3.0, 3.0, 2.3, 2.0],
        left=Inches(0.4), top=Inches(2.4), row_h_in=0.55)
    _fill_body_placeholder(_g2_slide, [
        (0, "Frame-by-frame agreement across 200 trials × 5 SNR levels:"),
        (0, " "),
        (0, " "),
        (0, " "),
        (0, "Conclusion: the Python decoder is correct. The QEX −1.4 dB reference reflects "
            "different test conditions — likely multi-frame SPD averaging or a different "
            "SNR convention. There is no decoder gap to close."),
        (1, "Remaining 1.4 dB gap is entirely in the SPD detection pipeline"),
    ])
    add_notes(_g2_slide,
        "This was the most important negative result of the investigation. "
        "Ruling out the decoder gap focused all remaining effort on the detection pipeline. "
        "Without the cross-language test it would have been tempting to chase phantom decoder bugs.")

    # Gap 3 ── SPD pipeline overhead
    content_slide(prs, "Gap 3: Oracle Testing — Isolating Pipeline Stages", [
        (0, "Oracle methodology: bypass each stage in turn and re-measure"),
        (1, "oracle_sync: skip detection, feed the correct frame position directly"),
        (1, "oracle_decode: skip detection + sync, feed the correct aligned frame"),
        (0, "The SPD pipeline has three stages that could cause missed decodes:"),
        (1, "Detection — squared-spectrum sliding window; candidate selection"),
        (1, "Sync — correlate against known sync word pattern; threshold 1.3"),
        (1, "Decoder — BP + LDPC on the averaged frame"),
        (0, "Three failure categories per trial:"),
        (1, "det_miss: detection never found the correct window"),
        (1, "sync_miss: detection found window, sync gate rejected it"),
        (1, "oracle_decode_miss: correct window + correct alignment, decoder still failed"),
        (0, "This cleanly attributes each missed decode to exactly one stage"),
    ], notes="The oracle concept is standard in algorithm analysis but rarely applied "
             "in amateur radio DSP work. It made it possible to answer 'where is the 1.4 dB?' "
             "definitively rather than through speculation.")

    # Gap 3 failure mode table
    _g3_slide = _blank(prs)
    add_rect(_g3_slide, 0, 0, W, Pt(8), ACCENT)
    add_textbox(_g3_slide, "Gap 3: SPD Failure Mode Profile (IIR path, 300 trials/SNR)",
                Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.75),
                color=NAVY, size=26, bold=True)
    _add_data_table(_g3_slide,
        headers=["SNR (dB)", "SPD%", "det_hit%", "det_metric", "oracle_sync%", "oracle_dec%", "det_miss", "sync_miss", "dec_miss"],
        rows=[
            ["-1.0",  "1.3%",  "9.3%",  "2.16", "47.7%", "5.0%",  "270", "15", "11"],
            ["-0.5",  "2.3%", "18.7%",  "2.39", "59.0%", "9.7%",  "241", "23", "29"],
            ["+0.0",  "8.0%", "31.3%",  "2.72", "74.7%","20.3%",  "202", "21", "53"],
            ["+0.5", "19.3%", "52.0%",  "3.16", "84.0%","33.3%",  "140", "24", "78"],
            ["+1.0", "36.3%", "67.7%",  "3.73", "91.7%","46.0%",   "92", "14", "85"],
            ["+1.5", "57.3%", "85.0%",  "4.47", "96.3%","58.7%",   "38",  "7", "83"],
            ["+2.0", "71.0%", "94.3%",  "5.42", "98.3%","66.7%",   "14",  "4", "69"],
            ["+2.5", "81.3%", "98.3%",  "6.60", "99.7%","73.3%",    "3",  "0", "53"],
        ],
        col_widths_in=[1.3, 0.9, 1.1, 1.2, 1.3, 1.3, 1.1, 1.1, 1.1],
        left=Inches(0.35), top=Inches(1.05), row_h_in=0.43, base_font_pt=12)
    add_textbox(_g3_slide,
        "At the +1.5 dB threshold: det_miss = 63.7% of failures, oracle_dec_miss = 29.4%.  "
        "Detection is the dominant bottleneck — not the decoder, not sync.",
        Inches(0.4), Inches(5.80), Inches(12.5), Inches(0.7),
        color=MUTED, size=16, italic=True)
    add_notes(_g3_slide,
        "det_miss dominates at threshold — the detector simply never finds the right window. "
        "oracle_dec at 58.7% vs SPD at 57.3% at +1.5 dB means detection overhead is only ~0.2 dB. "
        "The 1.4 dB SPD overhead is intrinsic to the algorithm, not fixable by tuning thresholds.")

    # BPF idea
    content_slide(prs, "The Audio BPF Insight", [
        (0, "Root cause of det_miss: noise bandwidth is too wide"),
        (1, "IIR decimation from 48→12 kHz passes noise up to ~4800 Hz"),
        (1, "Squared-spectrum detector integrates noise×noise cross-terms across that band"),
        (1, "At threshold SNR, these cross-terms dominate the detection metric"),
        (0, "Real SSB receivers band-limit audio to ~300–2700 Hz"),
        (1, "MSK144 signal occupies 1000–2000 Hz — well within that range"),
        (1, "Noise above 2700 Hz contributes nothing to the signal, only to the noise floor"),
        (0, "Hypothesis: model the SSB audio filter with a 300–2700 Hz BPF"),
        (1, "Noise bandwidth drops from ~4800 Hz to ~2400 Hz — a factor of 2"),
        (1, "In the squared domain, noise×noise cross-terms drop as BW²  — a factor of 4"),
        (1, "Predicted: significant improvement in detection metric"),
    ], notes="The SSB receiver analogy is the key insight. Real hardware operators "
             "hear through a 2.7 kHz audio passband — MAP144 was using 4.8 kHz of noise bandwidth. "
             "The BPF models what a receiver already does in hardware.")

    # BPF implementation
    content_slide(prs, "BPF Implementation — Getting the Math Right", [
        (0, "Naive approach: mix up → real BPF → Hilbert → mix down"),
        (1, "Re-introduces the Hilbert correlation problem from Gap 1"),
        (1, "Hilbert of filtered noise makes I and Q correlated — wrong LLRs in BP decoder"),
        (1, "Oracle decode rate: 0/8 at SNR=2 — complete failure"),
        (0, "Correct approach: complex lowpass applied independently to I and Q"),
        (1, "Signal at 0 Hz baseband; 300–2700 Hz BPF → symmetric ±1200 Hz lowpass"),
        (1, "Apply same real FIR to Re{z} and Im{z} separately, recombine"),
        (1, "I/Q statistical independence preserved — BP decoder works correctly"),
        (0, "Zero-phase filter (filtfilt) — no group delay, no onset transient"),
        (1, "Causal FIR at 12 kHz: 129-tap filter → 64-sample delay → corrupts 7% of frame"),
        (1, "filtfilt: zero phase, transients cancel at edges — clean 72 ms frame"),
    ], notes="Two implementation bugs were tried and diagnosed before the correct version. "
             "The first (Hilbert approach) looked correct but produced zero oracle decodes. "
             "The second (causal FIR) improved detection but hurt decoding. "
             "Only the filtfilt complex-lowpass on independent I/Q channels worked.")

    # BPF results comparison table
    _bpf_slide = _blank(prs)
    add_rect(_bpf_slide, 0, 0, W, Pt(8), ACCENT)
    add_textbox(_bpf_slide, "BPF Results: Detection Failures Eliminated (100 trials/SNR)",
                Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.75),
                color=NAVY, size=26, bold=True)
    _add_data_table(_bpf_slide,
        headers=["SNR (dB)", "No-BPF SPD%", "No-BPF det_hit%", "BPF-2700 SPD%", "BPF-2700 det_hit%", "BPF oracle_dec%"],
        rows=[
            ["-1.0",  "3%",  "9%",  "12%",  "100%",  "8%"],
            ["-0.5",  "4%", "20%",  "25%",  "100%", "17%"],
            ["+0.0",  "7%", "35%",  "42%",  "100%", "28%"],
            ["+0.5", "16%", "53%",  "55%",  "100%", "38%"],
            ["+1.0", "33%", "65%",  "68%",  "100%", "47%"],
            ["+1.5", "56%", "84%",  "83%",  "100%", "61%"],
            ["+2.0", "72%", "92%",  "89%",  "100%", "69%"],
            ["+2.5", "83%", "97%",  "91%",  "100%", "74%"],
        ],
        col_widths_in=[1.4, 1.7, 1.9, 1.7, 1.9, 1.9],
        left=Inches(0.35), top=Inches(1.05), row_h_in=0.43, base_font_pt=13)
    add_textbox(_bpf_slide,
        "With BPF: det_hit = 100% at every SNR level.  All remaining failures are oracle_decode_miss — "
        "the decoder limit.  The SPD is now decoder-limited, not detection-limited.",
        Inches(0.4), Inches(5.80), Inches(12.5), Inches(0.7),
        color=MUTED, size=16, italic=True)
    add_notes(_bpf_slide,
        "This is the most striking result in the investigation. "
        "Detection failures went from ~64% of all failures to zero. "
        "The detection metric jumped from 3–7 (no BPF) to 8–37 (BPF). "
        "The system is now limited by the decoder, not the front end.")

    # Final sensitivity budget
    _budget_slide = _title_content(prs)
    add_rect(_budget_slide, 0, 0, W, Pt(8), ACCENT)
    _fit_placeholders(_budget_slide)
    _set_title_placeholder(_budget_slide, "Final Sensitivity Budget")
    _fill_body_placeholder(_budget_slide, [
        (0, "Each gap measured independently; total improvement 4 dB over baseline:"),
    ])
    _add_data_table(_budget_slide,
        headers=["Fix", "SPD 50% threshold", "Cumulative gain", "Status"],
        rows=[
            ["Baseline — real audio, Hilbert path",  "+4.0 dB", "—",       "Measured"],
            ["Gap 1 — complex IQ, √2 scale fix",     "+1.4 dB", "+2.6 dB", "Done ✓"],
            ["Gap 2 — Python decoder accuracy",       "+1.4 dB", "0 dB",    "No gap exists ✓"],
            ["Gap 3 — audio BPF 300–2700 Hz",        "~0.0 dB", "+1.4 dB", "Done ✓"],
            ["Oracle decoder limit (floor)",         "~0.0 dB", "+4.0 dB", "Decoder-limited"],
            ["QEX reference (conditions unknown)",   "−1.4 dB", "—",       "Different conditions"],
        ],
        col_widths_in=[5.0, 2.8, 2.5, 2.2],
        left=Inches(0.4), top=Inches(2.2), row_h_in=0.50)
    add_notes(_budget_slide,
        "The QEX reference likely reflects the Fortran SPD with true multi-frame averaging "
        "across a real burst, not a single synthesized frame. "
        "Our 0.0 dB threshold matches the fundamental single-frame decoder limit exactly. "
        "Total gain from the investigation: 4 dB on short meteor pings.")

    content_slide(prs, "What Made This Investigation Possible", [
        (0, "Multi-session memory — findings persisted across days of work"),
        (1, "Intermediate results, failure modes, and rejected hypotheses all recorded"),
        (1, "New sessions started with full context, not re-deriving basics"),
        (0, "Claude executed complex tool chains without supervision"),
        (1, "Built Fortran test executables from source + wrote Python bridge"),
        (1, "Designed Monte Carlo test harness with oracle bypass modes"),
        (1, "Ran 300-trial sweeps, parsed results, updated conclusions"),
        (0, "Systematic decomposition prevented chasing phantoms"),
        (1, "Each gap tested and closed (or ruled out) before moving to the next"),
        (1, "Without Gap 2 verification, we might have chased a non-existent decoder bug"),
        (0, "Domain knowledge directed the investigation; AI executed it"),
        (1, "The BPF insight came from knowing how real SSB receivers work"),
        (1, "Claude implemented it; the engineer recognized why it should work"),
    ], notes="This is the meta-lesson: what distinguished this from a typical AI coding session "
             "was the systematic methodology. Each step was a falsifiable test with a clear result. "
             "The combination of human domain knowledge and AI execution speed is the leverage.")

    # ═══════════════════════════════════════════════════════════════════════════
    # SECTION 11 — MAP144 vs WSJT-X REAL-WORLD COMPARISON
    # ═══════════════════════════════════════════════════════════════════════════
    section_slide(prs, 11, "MAP144 vs WSJT-X",
        subtitle="Real-world overnight comparison",
        audience="Technical / operators",
        notes="Run MAP144 and WSJT-X simultaneously overnight on the same band. "
              "compare_decoders.py reads both logs in the morning and produces "
              "the comparison tables and plots.")

    content_slide(prs, "The Comparison Method", [
        (0, "MAP144 and WSJT-X run simultaneously on the same radio"),
        (1, "MAP144: FlexRadio DAX IQ — direct complex IQ at 48 kHz"),
        (1, "WSJT-X: FlexRadio DAX audio — standard SSB audio as digital sound in the PC"),
        (1, "Independent paths — no shared samples, no interference"),
        (0, "MAP144 logs every decode to MSK144/detections/decodes.jsonl"),
        (1, "One JSON line per decode: message, SNR, frequency, timestamp, decoder used"),
        (0, "WSJT-X logs every decode to ALL.TXT"),
        (1, "Standard WSJT-X format: date/time, freq, SNR, dt, audio-Hz, message"),
        (0, "compare_decoders.py reads both logs, matches on message + time (±1 s)"),
        (1, "Three categories: Both decoded  |  WSJT-X only  |  MAP144 only"),
        (0, "Deal with different multiple decodes per 15 second period"),
    ], notes="The independent radio paths are important — if they shared audio, "
             "any correlated noise would make the comparison misleading.")

    content_slide(prs, "compare_decoders.py — Output", [
        (0, "Per-ping classification table"),
        (1, "Both decoded / WSJT-X only / MAP144 only — with timestamp and frequency"),
        (1, "WSJT-X-only misses tagged: NOT_DETECTED / NO_DECODE / PERIOD_SLIP"),
        (0, "Coverage statistics"),
        (1, "WSJT-X coverage % and MAP144 coverage % of all unique pings observed"),
        (0, "SNR comparison on matched pairs"),
        (1, "MAP144 SNR field: est_snr_db (1500 ms pre-burst PSD, preferred)"),
        (1, "  jt9_snr_db on a short clip is unreliable — clusters at 0 dB"),
        (1, "Measured bias: WSJT-X SNR ≈ +1.7 dB over MAP144 est_snr_db"),
        (1, "  WSJT-X: coherent matched filter over full 15 s frame"),
        (1, "  MAP144: non-coherent peak-power vs pre-burst noise floor"),
        (0, "Optional plots: detection-rate curve and SNR scatter"),
        (1, "Detection-rate: MAP144 decode% vs WSJT-X SNR bin"),
        (1, "Scatter: (WSJT-X SNR, MAP144 est_snr) for matched pairs + stacked miss rail"),
    ], mono_items={},
    notes="The +1.7 dB bias (measured April 2026 overnight run) replaces the earlier "
          "~3–4 dB estimate that was inflated by jt9 short-clip SNR artifacts — "
          "jt9 has no noise baseline on a burst WAV and reported 0 dB for ~49% of decodes. "
          "est_snr_db uses the 1500 ms pre-burst window (jt9 fallback path) or "
          "500 ms (SPD path) and is the reliable SNR source. "
          "The detection-rate curve is the most useful single plot: "
          "it shows MAP144 sensitivity relative to WSJT-X as a function of SNR.")

    content_slide(prs, "Running the Morning Report", [
        (0, "After an overnight session:"),
        (1, "python compare_decoders.py --date YYYYMMDD --plot detect.png"),
        (0, "Produces:"),
        (1, "Console table: Both / WSJT-X only / MAP144 only"),
        (1, "detect.png — detection-rate curve by WSJT-X SNR bin"),
        (1, "detect_scatter.png — per-ping SNR scatter plot"),
        (0, "Quick launch tally from launches.jsonl:"),
        (1, "python -c \"import json; from collections import Counter; "
            "c = Counter(json.loads(l)['outcome'] for l in open('MSK144/detections/launches.jsonl') if l.strip()); print(c)\""),
        (0, "Healthy overnight run:"),
        (1, "Mostly no_decode (noise triggers) + some decoded + few timeout/error"),
        (1, "decoded / (decoded + no_decode) ratio tracks effective false alarm rate"),
    ], mono_items={1, 3, 4, 5, 7},
    notes="The launch tally is the fastest sanity check. "
          "If timeout count is high, jt9 is under CPU pressure. "
          "If no_decode >> decoded, the threshold may be too low.")

    # ═══════════════════════════════════════════════════════════════════════════
    # SECTION 12 — DOPPLER MEASUREMENTS & DESIGN OPPORTUNITIES
    # ═══════════════════════════════════════════════════════════════════════════
    section_slide(prs, 12, "Doppler Measurements & Design Opportunities",
        subtitle="What measured signal physics tells us about future decoder design",
        audience="Technical",
        notes="This section is the original research half of the project — measurements "
              "of within-burst Doppler statistics from a captured corpus, and the "
              "architectural extensions to WSJT-X that the measurements motivate.  "
              "Skip for a general-audience talk.")

    content_slide(prs, "Measuring Frequency vs Time", [
        (0, "Standard MS decoders (WSJT-X, jt9) report one frequency per decode"),
        (1, "No information about how the frequency changes during the burst"),
        (0, "MAP144 tool:  tools/measure_ping_freq_vs_time.py"),
        (1, "For each WSJT-X-saved WAV in a session corpus:"),
        (2, "Bandpass to MSK144 audio band (1300-1700 Hz)"),
        (2, "Square the signal — modulation tones become CW lines at 2·f₀ ± 1000 Hz"),
        (2, "Sliding short FFT, ±300 Hz peak search around each squared tone"),
        (2, "Carrier estimate = (low_tone + high_tone) / 4   at 10-ms resolution"),
        (1, "Suppresses MSK modulation contamination that breaks Hilbert-derivative methods"),
        (0, "Per-burst statistics computed:"),
        (1, "σ(f), Δf  —  frequency stability within the burst"),
        (1, "Pearson r(|A|, f)  —  direct test of the multipath hypothesis"),
        (1, "Burst duration, peak above noise floor"),
        (0, "These distributions are not in the published literature for 6 m MS"),
    ], mono_items={2, 3},
    notes="The squared-spectrum approach is the same trick MAP144's tone-pair detector uses. "
          "Hilbert-phase-derivative measurements gave inflated σ(f) due to MSK modulation leak — "
          "squared FFT is cleaner.  Validated against the same WAV with both methods.")

    image_slide(prs, "Ping vs Continuous Signal — Same Metric, Different Physics",
        label="doppler_ping_vs_continuous",
        caption="Left: meteor-scatter ping (KB8OTK, 90 ms, σ_f = 7.7 Hz).  "
                "Right: strong-local continuous signal (N1KWF, 15 s, σ_f = 44.8 Hz).",
        notes="The ping has stable frequency over its short lifetime — coherent integration "
              "could in principle gain on it.  The continuous signal drifts substantially "
              "over 15 s — a single frequency estimate is wrong by tens of Hz at the "
              "session edges.  Period-mode decode (TODO #42) needs in-window tracking.",
        png="doppler_ping_vs_continuous.png")

    image_slide(prs, "Doppler Statistics from 145 Bursts",
        label="doppler_burst_stats_distributions",
        caption="2026-05-11 corpus.  Top-left: σ(f) histogram with coherent-gain ceilings.  "
                "Top-right: σ(f) vs duration.  Bottom-left: |r(|A|, f)| distribution.  "
                "Bottom-right: |r| vs σ(f), peak-SNR coloured.",
        notes="Vertical red dashed lines in top-left are the coherent-gain ceilings for "
              "N=1,2,3,5,7 frames.  Most bursts exceed even the 1-frame ceiling.  "
              "Bottom-right shows that bursts with higher σ(f) tend to have stronger |r| — "
              "the multipath signature.",
        png="doppler_burst_stats_distributions.png")

    content_slide(prs, "Findings from 145 Bursts", [
        (0, "σ(f) within burst"),
        (1, "median 7.5 Hz, p75 13.1 Hz, max 128 Hz"),
        (1, "Grows monotonically with burst duration"),
        (2, "<100 ms bursts:  median 6.6 Hz"),
        (2, "100-200 ms:        median 10.8 Hz"),
        (2, "200-400 ms:        median 13.5 Hz"),
        (0, "Fraction satisfying the worst-case coherent-gain ceiling"),
        (1, "N=1 frame  (σ<1.70 Hz):  4.8% of bursts"),
        (1, "N=2 frames (σ<0.87 Hz):  1.4%"),
        (1, "N=3 frames (σ<0.58 Hz):  0%   ← zero"),
        (1, "N=5, N=7:                0%   ← zero"),
        (0, "Amplitude-frequency correlation"),
        (1, "Strong (|r|>0.5):    26% of bursts"),
        (1, "Moderate (|r|>0.3):  50%"),
        (1, "Confirms the multipath mechanism — flutter is correlated"),
        (0, "Ceiling formula is conservative (worst-case random walk)"),
        (1, "Linear drift is easier — decoder partially tolerates it"),
        (1, "But the trend is unambiguous: 5/7-frame integration is Doppler-limited"),
    ], notes="The N=3 row at 0% is the striking result.  Conservative formula, yes, "
             "but the message holds: -d 3 Deep mode's longer averages have very little "
             "physics-allowed signal to integrate on real-world pings at this site.")

    content_slide(prs, "Architectural Opportunities for MAP144", [
        (0, "WSJT-X averages with no in-burst frequency tracking"),
        (1, "Coherent gain degrades silently when signal drifts"),
        (1, "Operator sees no signal that drift caused the miss"),
        (0, "MAP144 has the data to do better"),
        (1, "Already landed: detection-stage fine-frequency pre-correction"),
        (2, "Sync correlator → audio mix shift → jt9 FTOL drops 200 → 50 Hz"),
        (2, "~4× CPU reduction per decode attempt"),
        (1, "Next step: per-burst drift correction inside the integration loop"),
        (2, "Split burst into sub-windows, estimate f(t) per window"),
        (2, "Fit linear (or polynomial) drift model"),
        (2, "Mix-correct against the drift  before  coherent averaging"),
        (2, "Recovers coherent gain on linearly drifting bursts (most pings)"),
        (0, "Gate the averaging length"),
        (1, "Measure σ(f) on-the-fly during sub-window scan"),
        (1, "Use N=2 if σ < 1 Hz, N=1 otherwise"),
        (1, "Stops decoder from trying integrations physics won't support"),
        (0, "Genuine algorithmic step beyond WSJT-X — measurement-driven, not a port"),
    ], notes="The per-burst drift correction is the key idea — WSJT-X assumes the signal is "
             "stationary, MAP144 would estimate the drift and undo it before averaging.  "
             "TODO #20 and TODO #42 are versions of this.")

    # ═══════════════════════════════════════════════════════════════════════════
    # SECTION 13 — SUPPORT
    # ═══════════════════════════════════════════════════════════════════════════
    section_slide(prs, 13, "Support and Resources",
        audience="Everyone",
        notes="Close on a light note.")

    content_slide(prs, "Source Code and Documentation", [
        (0, "GitHub repository:"),
        (1, "github.com/wa1hco/map144"),
        (0, "README covers installation and basic operation"),
        (0, "Issues tracker: bug reports and feature requests welcome"),
        (0, "WAV test files included for offline testing"),
        (0, "requirements.txt lists all Python dependencies"),
    ], mono_items={1},
    notes="Replace with actual GitHub URL.")

    quote_slide(prs,
        "If you have a question about the code…\n"
        "don't ask me.\n"
        "Ask Claude.",
        attribution="WA1HCO, 2026",
        notes="Honest answer: the AI that helped write it can also explain it. "
              "This usually gets a laugh. Follow up with: 'Claude has read every line of this code'.")

    content_slide(prs, "Getting Help", [
        (0, "For installation and operation questions:"),
        (1, "This is 2026, just ask AI and/or feed it error messages"),
        (1, "AI knows this stuff _way_ better than I"),
        (0, "For code questions:"),
        (1, "Ask Claude (claude.ai) — point to or paste the relevant code and ask for explanation"),
        (1, "Claude wrote it and can explain any of it"),
        (0, "For MSK144 / meteor scatter questions:"),
        (1, "Claude knows all the published literature"),
        (1, "WSJT-X community: wsjt.sourceforge.io"),
        (1, "ON4KST chat during meteor showers"),
        (1, "Your local VHF/UHF club"),
        (0, "For propagation data: pskreporter.info"),
    ], notes="Redirect technical questions to GitHub issues so answers are searchable.")

    content_slide(prs, "Thank You", [
        (0, "Jeff Millar  WA1HCO"),
        (0, " "),
        (0, "Questions?"),
        (0, " "),
        (0, "github.com/wa1hco/map144"),
    ], notes="Leave time for questions. Common ones:\n"
             "- Does it work on Windows? (Mostly, with WSL)\n"
             "- Can it decode FT8? (No — different protocol, different tool)\n"
             "- What antenna do you use? (Your favourite question)")

    # ── Footers ───────────────────────────────────────────────────────────────
    for _i, _sl in enumerate(prs.slides, start=1):
        _add_footer(_sl, _i)

    # ── Save ──────────────────────────────────────────────────────────────────
    out_path = Path(__file__).parent / "map144_presentation.pptx"
    prs.save(str(out_path))
    print(f"Saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
